"""Format-aware validation for generated artifacts before delivery.

The artifact store intentionally remains format-agnostic.  This module owns
the narrower contract that a generated ``.pptx`` must be an objectively
readable OOXML presentation before a delivery surface can call it published.
"""

from __future__ import annotations

import io
import posixpath
import struct
import xml.etree.ElementTree as ET
import xml.parsers.expat as expat
import zipfile
import zlib
from dataclasses import dataclass
from typing import Never
from urllib.parse import urlsplit

import structlog

from opensquilla.contracts.attachments import OLE_MAGIC, PPTX_MIME

log = structlog.get_logger(__name__)

__all__ = [
    "PPTX_MIME",
    "ArtifactValidationError",
    "ArtifactValidationReport",
    "is_pptx_candidate",
    "validate_pptx_bytes",
    "validate_artifact_for_delivery",
]

_MAX_INFLATED_BYTES = 200 * 1024 * 1024
_MAX_CORE_XML_BYTES = 16 * 1024 * 1024
_MAX_ZIP_MEMBERS = 10_000
_MAX_CENTRAL_DIRECTORY_BYTES = 32 * 1024 * 1024
_READ_CHUNK_BYTES = 1024 * 1024

_LOCAL_FILE_HEADER = b"PK\x03\x04"
_CENTRAL_DIRECTORY_HEADER = b"PK\x01\x02"
_END_OF_CENTRAL_DIRECTORY = b"PK\x05\x06"
_ZIP64_END_OF_CENTRAL_DIRECTORY = b"PK\x06\x06"
_ZIP64_END_OF_CENTRAL_DIRECTORY_LOCATOR = b"PK\x06\x07"
_MAX_ZIP_COMMENT_BYTES = 65_535
_EOCD_BYTES = 22
_LOCAL_FILE_HEADER_BYTES = 30
_CENTRAL_DIRECTORY_HEADER_BYTES = 46
_ZIP64_EOCD_MIN_SIZE = 44
_ZIP64_EOCD_LOCATOR_BYTES = 20
_DATA_DESCRIPTOR_SIGNATURE = b"PK\x07\x08"
_ZIP64_EXTRA_FIELD_ID = 0x0001
_SUPPORTED_ZIP_COMPRESSION = frozenset({zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED})
_ASCII_HEX_DIGITS = frozenset("0123456789abcdefABCDEF")
_ASCII_UNRESERVED_BYTES = frozenset(
    b"ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789-._~"
)
_ASCII_LOWER_TRANSLATION = str.maketrans(
    "ABCDEFGHIJKLMNOPQRSTUVWXYZ",
    "abcdefghijklmnopqrstuvwxyz",
)

_CONTENT_TYPES_PART = "[Content_Types].xml"
_ROOT_RELS_PART = "_rels/.rels"

_CONTENT_TYPES_NAMESPACES = frozenset(
    {
        "http://schemas.openxmlformats.org/package/2006/content-types",
        "http://purl.oclc.org/ooxml/package/content-types",
    }
)
_PACKAGE_RELATIONSHIP_NAMESPACES = frozenset(
    {
        "http://schemas.openxmlformats.org/package/2006/relationships",
        "http://purl.oclc.org/ooxml/package/relationships",
    }
)

_TRANSITIONAL_RELATIONSHIP_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
_STRICT_RELATIONSHIP_NS = "http://purl.oclc.org/ooxml/officeDocument/relationships"
_OFFICE_DOCUMENT_RELATIONSHIP_TYPES = frozenset(
    {
        f"{_TRANSITIONAL_RELATIONSHIP_NS}/officeDocument",
        f"{_STRICT_RELATIONSHIP_NS}/officeDocument",
    }
)

_TRANSITIONAL_PRESENTATION_NS = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)
_STRICT_PRESENTATION_NS = "http://purl.oclc.org/ooxml/presentationml/main"

_PRESENTATION_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_SLIDE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)

_GENERIC_USER_MESSAGE = (
    "The PPTX was not attached because it is not a readable OOXML presentation. "
    "Repair or regenerate it and try again."
)
_ENCRYPTED_USER_MESSAGE = (
    "The PPTX was not attached because encrypted or legacy Office containers are not "
    "supported. Save it as an unencrypted .pptx, or regenerate it, then try again."
)
_INFLATED_LIMIT_USER_MESSAGE = (
    "The PPTX was not attached because its expanded contents exceed the safe validation "
    "limit. Reduce or regenerate it and try again."
)
_PACKAGE_COMPLEXITY_USER_MESSAGE = (
    "The PPTX was not attached because its package structure exceeds the safe validation "
    "limit. Simplify or regenerate it and try again."
)


class ArtifactValidationError(ValueError):
    """Blocking artifact validation failure safe to expose to a tool caller."""

    def __init__(self, reason_code: str, user_message: str) -> None:
        super().__init__(user_message)
        self.reason_code = reason_code
        self.user_message = user_message


class _UnsafeXmlDeclarationError(ValueError):
    """Raised by the non-resolving XML preflight when a DTD or entity is declared."""


@dataclass(frozen=True)
class ArtifactValidationReport:
    """Successful validation result, optionally carrying compatibility warnings."""

    warnings: tuple[str, ...] = ()


@dataclass(frozen=True)
class _Relationship:
    relationship_id: str
    relationship_type: str
    target: str
    target_mode: str | None


@dataclass(frozen=True)
class _ContentTypes:
    defaults: dict[str, str]
    overrides: dict[str, str]

    def for_part(self, part_name: str) -> str | None:
        override = self.overrides.get(_canonical_part_key(part_name))
        if override is not None:
            return override
        filename = posixpath.basename(part_name)
        extension = filename.rsplit(".", 1)[1].casefold() if "." in filename else ""
        return self.defaults.get(extension)


@dataclass(frozen=True)
class _CentralDirectoryEntry:
    decoded_name: str
    raw_name: bytes
    version_needed: int
    flags: int
    compression: int
    modified_time: int
    modified_date: int
    crc: int
    compressed_size: int
    uncompressed_size: int
    local_header_offset: int
    uses_zip64_sizes: bool


@dataclass(frozen=True)
class _PackageParts:
    """Case-insensitive OPC part keys mapped to their physical ZIP member names."""

    by_key: dict[str, str]
    case_fallback_names: dict[str, str]

    def contains(self, part_name: str) -> bool:
        return _canonical_part_key(part_name) in self.by_key

    def physical_name(self, part_name: str) -> str:
        physical_name = self.by_key.get(_canonical_part_key(part_name))
        if physical_name is None:
            raise KeyError(part_name)
        return physical_name

    def reference_requires_case_fallback(self, part_name: str) -> bool:
        physical_name = self.physical_name(part_name)
        normalised_reference = _normalise_package_path(part_name, base_parts=[])
        normalised_physical = _normalise_package_path(physical_name, base_parts=[])
        requires_fallback = normalised_reference != normalised_physical
        if requires_fallback:
            self.case_fallback_names.setdefault(physical_name, normalised_reference)
        return requires_fallback


def is_pptx_candidate(
    *,
    source_name: str | None = None,
    name: str | None = None,
    mime: str | None = None,
) -> bool:
    """Return whether any delivery signal identifies the artifact as PPTX."""

    names = (source_name, name)
    if any(
        isinstance(value, str) and value.strip().casefold().endswith(".pptx")
        for value in names
    ):
        return True
    if not isinstance(mime, str):
        return False
    return mime.split(";", 1)[0].strip().casefold() == PPTX_MIME


def validate_artifact_for_delivery(
    payload: bytes,
    *,
    source_name: str | None,
    name: str,
    mime: str,
    source: str,
) -> ArtifactValidationReport:
    """Validate a generated artifact at the final format-aware delivery boundary.

    Non-PPTX artifacts intentionally remain opaque.  Logs contain only bounded
    operational metadata; names, paths, exception strings, and content are not
    emitted.
    """

    if not is_pptx_candidate(source_name=source_name, name=name, mime=mime):
        return ArtifactValidationReport()

    try:
        report = validate_pptx_bytes(payload)
    except ArtifactValidationError as exc:
        log.warning(
            "artifact_validation.completed",
            source=source,
            outcome="blocked",
            reason_code=exc.reason_code,
            size=len(payload),
        )
        raise

    if report.warnings:
        for warning in report.warnings:
            log.warning(
                "artifact_validation.completed",
                source=source,
                outcome="warning",
                reason_code=warning,
                size=len(payload),
            )
    else:
        log.info(
            "artifact_validation.completed",
            source=source,
            outcome="passed",
            reason_code="ok",
            size=len(payload),
        )
    return report


def validate_pptx_bytes(payload: bytes) -> ArtifactValidationReport:
    """Validate PPTX bytes and return any non-blocking compatibility warnings."""

    if payload.startswith(OLE_MAGIC):
        _fail("pptx_encrypted_or_legacy_container", _ENCRYPTED_USER_MESSAGE)

    central_entries = _validate_zip_envelope(payload)
    try:
        with zipfile.ZipFile(io.BytesIO(payload), "r") as archive:
            infos = archive.infolist()
            if len(infos) != len(central_entries):
                _fail("pptx_invalid_central_directory")

            parts = _validate_zip_members(archive, infos, central_entries)
            content_types = _parse_content_types(archive, parts)
            root_relationships = _parse_relationships(
                archive,
                parts,
                _ROOT_RELS_PART,
                missing_reason="pptx_missing_root_relationships",
            )
            content_types_case_fallback = parts.reference_requires_case_fallback(
                _CONTENT_TYPES_PART
            )
            root_rels_case_fallback = parts.reference_requires_case_fallback(_ROOT_RELS_PART)
            parser_case_compatibility = (
                content_types_case_fallback or root_rels_case_fallback
            )
            office_relationships = [
                relationship
                for relationship in root_relationships.values()
                if relationship.relationship_type in _OFFICE_DOCUMENT_RELATIONSHIP_TYPES
            ]
            if len(office_relationships) != 1:
                _fail("pptx_invalid_office_document_relationship")

            office_relationship = office_relationships[0]
            if office_relationship.target_mode == "External":
                _fail("pptx_external_office_document")
            presentation_part = _resolve_relationship_target(
                source_part=None,
                target=office_relationship.target,
            )
            if not parts.contains(presentation_part):
                _fail("pptx_missing_presentation_part")
            presentation_case_fallback = parts.reference_requires_case_fallback(
                presentation_part
            )
            parser_case_compatibility = (
                presentation_case_fallback or parser_case_compatibility
            )
            presentation_content_type = content_types.for_part(presentation_part)
            if (
                presentation_content_type is None
                or presentation_content_type.casefold()
                != _PRESENTATION_CONTENT_TYPE.casefold()
            ):
                _fail("pptx_invalid_presentation_content_type")
            parser_case_compatibility = (
                parser_case_compatibility
                or presentation_content_type != _PRESENTATION_CONTENT_TYPE
            )

            is_strict = office_relationship.relationship_type.startswith(
                _STRICT_RELATIONSHIP_NS
            )
            parser_case_compatibility = _validate_presentation(
                archive,
                parts,
                content_types,
                presentation_part,
                is_strict=is_strict,
            ) or parser_case_compatibility
            parser_case_compatibility = (
                _validate_reachable_relationship_targets(
                    archive,
                    parts,
                    content_types,
                    root_relationships,
                    presentation_part=presentation_part,
                )
                or parser_case_compatibility
            )
    except ArtifactValidationError:
        raise
    except NotImplementedError:
        _fail("pptx_unsupported_zip_feature")
    except (
        zipfile.BadZipFile,
        zipfile.LargeZipFile,
        zlib.error,
        EOFError,
        OSError,
        RuntimeError,
    ):
        _fail("pptx_corrupt_zip_member")
    except ValueError:
        _fail("pptx_unsupported_zip_feature")

    if is_strict:
        return ArtifactValidationReport(warnings=("pptx_strict_ooxml_smoke_skipped",))
    return _python_pptx_smoke_test(
        payload,
        allow_opc_case_compatibility=parser_case_compatibility,
        case_fallback_names=parts.case_fallback_names,
    )


def _validate_zip_envelope(payload: bytes) -> tuple[_CentralDirectoryEntry, ...]:
    if not payload.startswith(_LOCAL_FILE_HEADER):
        _fail("pptx_not_zip")

    minimum_offset = max(0, len(payload) - (_EOCD_BYTES + _MAX_ZIP_COMMENT_BYTES))
    search_end = len(payload)
    eocd_offset = -1
    eocd_fields: tuple[int, ...] | None = None
    while search_end > minimum_offset:
        candidate = payload.rfind(_END_OF_CENTRAL_DIRECTORY, minimum_offset, search_end)
        if candidate < 0:
            break
        if candidate + _EOCD_BYTES <= len(payload):
            try:
                fields = struct.unpack_from("<4H2LH", payload, candidate + 4)
            except struct.error:
                _fail("pptx_invalid_zip_envelope")
            comment_length = fields[-1]
            if candidate + _EOCD_BYTES + comment_length == len(payload):
                eocd_offset = candidate
                eocd_fields = fields
                break
        search_end = candidate

    if eocd_fields is None:
        _fail("pptx_invalid_zip_envelope")

    disk_number, central_disk, disk_entries, total_entries, size, offset, _ = eocd_fields
    uses_zip64 = (
        disk_number == 0xFFFF
        or central_disk == 0xFFFF
        or disk_entries == 0xFFFF
        or total_entries == 0xFFFF
        or size == 0xFFFFFFFF
        or offset == 0xFFFFFFFF
    )
    central_directory_end = eocd_offset
    if uses_zip64:
        (
            disk_number,
            central_disk,
            disk_entries,
            total_entries,
            size,
            offset,
            central_directory_end,
        ) = _read_zip64_directory_fields(
            payload,
            eocd_offset=eocd_offset,
            legacy_fields=eocd_fields,
        )

    if disk_number != 0 or central_disk != 0 or disk_entries != total_entries:
        _fail("pptx_multidisk_zip_unsupported")
    if total_entries > _MAX_ZIP_MEMBERS:
        _fail("pptx_member_count_limit", _PACKAGE_COMPLEXITY_USER_MESSAGE)
    if size > _MAX_CENTRAL_DIRECTORY_BYTES:
        _fail("pptx_central_directory_limit", _PACKAGE_COMPLEXITY_USER_MESSAGE)
    if offset > central_directory_end or size != central_directory_end - offset:
        _fail("pptx_invalid_central_directory")
    entries = _validate_central_directory(
        payload,
        offset=offset,
        size=size,
        expected_entries=total_entries,
    )
    _validate_local_file_layout(payload, entries, central_directory_offset=offset)
    return entries


def _read_zip64_directory_fields(
    payload: bytes,
    *,
    eocd_offset: int,
    legacy_fields: tuple[int, ...],
) -> tuple[int, int, int, int, int, int, int]:
    locator_offset = eocd_offset - _ZIP64_EOCD_LOCATOR_BYTES
    if (
        locator_offset < 0
        or payload[locator_offset : locator_offset + 4]
        != _ZIP64_END_OF_CENTRAL_DIRECTORY_LOCATOR
    ):
        _fail("pptx_invalid_zip64_directory")
    try:
        locator_disk, zip64_offset, total_disks = struct.unpack_from(
            "<LQL", payload, locator_offset + 4
        )
    except struct.error:
        _fail("pptx_invalid_zip64_directory")
    if locator_disk != 0 or total_disks != 1:
        _fail("pptx_multidisk_zip_unsupported")
    if (
        zip64_offset < 0
        or zip64_offset + 12 > locator_offset
        or payload[zip64_offset : zip64_offset + 4] != _ZIP64_END_OF_CENTRAL_DIRECTORY
    ):
        _fail("pptx_invalid_zip64_directory")
    try:
        record_size = struct.unpack_from("<Q", payload, zip64_offset + 4)[0]
    except struct.error:
        _fail("pptx_invalid_zip64_directory")
    if (
        record_size < _ZIP64_EOCD_MIN_SIZE
        or zip64_offset + 12 + record_size != locator_offset
    ):
        _fail("pptx_invalid_zip64_directory")
    try:
        (
            _version_made_by,
            _version_needed,
            disk_number,
            central_disk,
            disk_entries,
            total_entries,
            size,
            offset,
        ) = struct.unpack_from("<2H2L4Q", payload, zip64_offset + 12)
    except struct.error:
        _fail("pptx_invalid_zip64_directory")

    (
        legacy_disk,
        legacy_central_disk,
        legacy_disk_entries,
        legacy_total_entries,
        legacy_size,
        legacy_offset,
        _,
    ) = legacy_fields
    pairs = (
        (legacy_disk, 0xFFFF, disk_number),
        (legacy_central_disk, 0xFFFF, central_disk),
        (legacy_disk_entries, 0xFFFF, disk_entries),
        (legacy_total_entries, 0xFFFF, total_entries),
        (legacy_size, 0xFFFFFFFF, size),
        (legacy_offset, 0xFFFFFFFF, offset),
    )
    if any(legacy != sentinel and legacy != actual for legacy, sentinel, actual in pairs):
        _fail("pptx_invalid_zip64_directory")
    return (
        disk_number,
        central_disk,
        disk_entries,
        total_entries,
        size,
        offset,
        zip64_offset,
    )


def _validate_central_directory(
    payload: bytes,
    *,
    offset: int,
    size: int,
    expected_entries: int,
) -> tuple[_CentralDirectoryEntry, ...]:
    cursor = offset
    end = offset + size
    entries: list[_CentralDirectoryEntry] = []
    while cursor < end:
        if len(entries) >= _MAX_ZIP_MEMBERS:
            _fail("pptx_member_count_limit", _PACKAGE_COMPLEXITY_USER_MESSAGE)
        if (
            cursor + _CENTRAL_DIRECTORY_HEADER_BYTES > end
            or payload[cursor : cursor + 4] != _CENTRAL_DIRECTORY_HEADER
        ):
            _fail("pptx_invalid_central_directory")
        try:
            (
                _version_made_by,
                version_needed,
                flags,
                compression,
                modified_time,
                modified_date,
                crc,
                raw_compressed_size,
                raw_uncompressed_size,
                filename_length,
                extra_length,
                comment_length,
                disk_start,
                _internal_attributes,
                _external_attributes,
                raw_local_header_offset,
            ) = struct.unpack_from(
                "<6H3L5H2L", payload, cursor + 4
            )
        except struct.error:
            _fail("pptx_invalid_central_directory")
        record_end = (
            cursor
            + _CENTRAL_DIRECTORY_HEADER_BYTES
            + filename_length
            + extra_length
            + comment_length
        )
        if record_end > end:
            _fail("pptx_invalid_central_directory")
        name_start = cursor + _CENTRAL_DIRECTORY_HEADER_BYTES
        raw_name = payload[name_start : name_start + filename_length]
        extra_start = name_start + filename_length
        extra = payload[extra_start : extra_start + extra_length]
        decoded_name = _decode_and_validate_zip_member_name(raw_name, flags=flags)
        (
            uncompressed_size,
            compressed_size,
            local_header_offset,
            resolved_disk_start,
        ) = _resolve_zip64_entry_fields(
            extra,
            uncompressed_size=raw_uncompressed_size,
            compressed_size=raw_compressed_size,
            local_header_offset=raw_local_header_offset,
            disk_start=disk_start,
            reason_code="pptx_invalid_central_directory",
        )
        if local_header_offset is None or resolved_disk_start is None:
            _fail("pptx_invalid_central_directory")
        if resolved_disk_start != 0:
            _fail("pptx_multidisk_zip_unsupported")
        if flags & 0x1:
            _fail("pptx_encrypted_zip_member", _ENCRYPTED_USER_MESSAGE)
        if compression not in _SUPPORTED_ZIP_COMPRESSION:
            _fail("pptx_unsupported_zip_compression")
        entries.append(
            _CentralDirectoryEntry(
                decoded_name=decoded_name,
                raw_name=raw_name,
                version_needed=version_needed,
                flags=flags,
                compression=compression,
                modified_time=modified_time,
                modified_date=modified_date,
                crc=crc,
                compressed_size=compressed_size,
                uncompressed_size=uncompressed_size,
                local_header_offset=local_header_offset,
                uses_zip64_sizes=(
                    raw_compressed_size == 0xFFFFFFFF
                    or raw_uncompressed_size == 0xFFFFFFFF
                ),
            )
        )
        cursor = record_end
    if cursor != end or len(entries) != expected_entries:
        _fail("pptx_invalid_central_directory")
    return tuple(entries)


def _resolve_zip64_entry_fields(
    extra: bytes,
    *,
    uncompressed_size: int,
    compressed_size: int,
    local_header_offset: int | None,
    disk_start: int | None,
    reason_code: str,
) -> tuple[int, int, int | None, int | None]:
    zip64_payload: bytes | None = None
    cursor = 0
    while cursor < len(extra):
        if cursor + 4 > len(extra):
            _fail(reason_code)
        field_id, field_size = struct.unpack_from("<HH", extra, cursor)
        cursor += 4
        field_end = cursor + field_size
        if field_end > len(extra):
            _fail(reason_code)
        if field_id == _ZIP64_EXTRA_FIELD_ID:
            if zip64_payload is not None:
                _fail(reason_code)
            zip64_payload = extra[cursor:field_end]
        cursor = field_end

    required = (
        uncompressed_size == 0xFFFFFFFF
        or compressed_size == 0xFFFFFFFF
        or local_header_offset == 0xFFFFFFFF
        or disk_start == 0xFFFF
    )
    if not required:
        return uncompressed_size, compressed_size, local_header_offset, disk_start
    if zip64_payload is None:
        _fail(reason_code)

    value_offset = 0

    def read_value(width: int) -> int:
        nonlocal value_offset
        if value_offset + width > len(zip64_payload):
            _fail(reason_code)
        format_code = "<Q" if width == 8 else "<L"
        value = int(struct.unpack_from(format_code, zip64_payload, value_offset)[0])
        value_offset += width
        return value

    if uncompressed_size == 0xFFFFFFFF:
        uncompressed_size = read_value(8)
    if compressed_size == 0xFFFFFFFF:
        compressed_size = read_value(8)
    if local_header_offset == 0xFFFFFFFF:
        local_header_offset = read_value(8)
    if disk_start == 0xFFFF:
        disk_start = read_value(4)
    return uncompressed_size, compressed_size, local_header_offset, disk_start


def _validate_local_file_layout(
    payload: bytes,
    entries: tuple[_CentralDirectoryEntry, ...],
    *,
    central_directory_offset: int,
) -> None:
    ordered = sorted(entries, key=lambda entry: entry.local_header_offset)
    if not ordered or ordered[0].local_header_offset != 0:
        _fail("pptx_invalid_local_file_layout")
    if len({entry.local_header_offset for entry in ordered}) != len(ordered):
        _fail("pptx_invalid_local_file_layout")

    for index, entry in enumerate(ordered):
        offset = entry.local_header_offset
        boundary = (
            ordered[index + 1].local_header_offset
            if index + 1 < len(ordered)
            else central_directory_offset
        )
        if (
            offset < 0
            or offset + _LOCAL_FILE_HEADER_BYTES > boundary
            or payload[offset : offset + 4] != _LOCAL_FILE_HEADER
        ):
            _fail("pptx_invalid_local_file_header")
        try:
            (
                version_needed,
                flags,
                compression,
                modified_time,
                modified_date,
                local_crc,
                raw_compressed_size,
                raw_uncompressed_size,
                filename_length,
                extra_length,
            ) = struct.unpack_from("<5H3L2H", payload, offset + 4)
        except struct.error:
            _fail("pptx_invalid_local_file_header")
        name_start = offset + _LOCAL_FILE_HEADER_BYTES
        data_start = name_start + filename_length + extra_length
        if data_start > boundary:
            _fail("pptx_invalid_local_file_header")
        raw_name = payload[name_start : name_start + filename_length]
        extra = payload[name_start + filename_length : data_start]
        _decode_and_validate_zip_member_name(raw_name, flags=flags)

        if (
            raw_name != entry.raw_name
            or version_needed != entry.version_needed
            or flags != entry.flags
            or compression != entry.compression
            or modified_time != entry.modified_time
            or modified_date != entry.modified_date
        ):
            _fail("pptx_invalid_local_file_header")
        if compression not in _SUPPORTED_ZIP_COMPRESSION:
            _fail("pptx_unsupported_zip_compression")

        uses_data_descriptor = bool(flags & 0x08)
        local_uses_zip64_sizes = (
            raw_compressed_size == 0xFFFFFFFF
            or raw_uncompressed_size == 0xFFFFFFFF
        )
        if uses_data_descriptor:
            (
                local_uncompressed_size,
                local_compressed_size,
                _,
                _,
            ) = _resolve_zip64_entry_fields(
                extra,
                uncompressed_size=raw_uncompressed_size,
                compressed_size=raw_compressed_size,
                local_header_offset=None,
                disk_start=None,
                reason_code="pptx_invalid_local_file_header",
            )
            if local_crc not in {0, entry.crc}:
                _fail("pptx_invalid_local_file_header")
            if local_compressed_size not in {0, entry.compressed_size}:
                _fail("pptx_invalid_local_file_header")
            if local_uncompressed_size not in {0, entry.uncompressed_size}:
                _fail("pptx_invalid_local_file_header")
        else:
            (
                local_uncompressed_size,
                local_compressed_size,
                _,
                _,
            ) = _resolve_zip64_entry_fields(
                extra,
                uncompressed_size=raw_uncompressed_size,
                compressed_size=raw_compressed_size,
                local_header_offset=None,
                disk_start=None,
                reason_code="pptx_invalid_local_file_header",
            )
            if (
                local_crc != entry.crc
                or local_compressed_size != entry.compressed_size
                or local_uncompressed_size != entry.uncompressed_size
            ):
                _fail("pptx_invalid_local_file_header")

        data_end = data_start + entry.compressed_size
        if data_end > boundary:
            _fail("pptx_invalid_local_file_layout")
        if uses_data_descriptor:
            _validate_data_descriptor(
                payload[data_end:boundary],
                entry,
                uses_zip64=entry.uses_zip64_sizes or local_uses_zip64_sizes,
            )
        elif data_end != boundary:
            _fail("pptx_invalid_local_file_layout")


def _validate_data_descriptor(
    descriptor: bytes,
    entry: _CentralDirectoryEntry,
    *,
    uses_zip64: bool,
) -> None:
    expected_payload_size = 20 if uses_zip64 else 12
    if len(descriptor) == expected_payload_size:
        signature_bytes = 0
    elif (
        len(descriptor) == expected_payload_size + 4
        and descriptor.startswith(_DATA_DESCRIPTOR_SIGNATURE)
    ):
        signature_bytes = 4
    else:
        _fail("pptx_invalid_data_descriptor")
    cursor = signature_bytes
    try:
        crc = struct.unpack_from("<L", descriptor, cursor)[0]
        cursor += 4
        if uses_zip64:
            compressed_size, uncompressed_size = struct.unpack_from("<QQ", descriptor, cursor)
        else:
            compressed_size, uncompressed_size = struct.unpack_from("<LL", descriptor, cursor)
    except struct.error:
        _fail("pptx_invalid_data_descriptor")
    if (
        crc != entry.crc
        or compressed_size != entry.compressed_size
        or uncompressed_size != entry.uncompressed_size
    ):
        _fail("pptx_invalid_data_descriptor")


def _validate_zip_members(
    archive: zipfile.ZipFile,
    infos: list[zipfile.ZipInfo],
    central_entries: tuple[_CentralDirectoryEntry, ...],
) -> _PackageParts:
    names: set[str] = set()
    parts_by_key: dict[str, str] = {}
    total = 0
    for info, central_entry in zip(infos, central_entries, strict=True):
        name = info.filename
        if (
            name != central_entry.decoded_name
            or info.header_offset != central_entry.local_header_offset
            or info.flag_bits != central_entry.flags
            or info.compress_type != central_entry.compression
            or info.CRC != central_entry.crc
            or info.compress_size != central_entry.compressed_size
            or info.file_size != central_entry.uncompressed_size
        ):
            _fail("pptx_invalid_central_directory")
        if name in names:
            _fail("pptx_duplicate_zip_member")
        _validate_member_name(name)
        names.add(name)
        part_key = _canonical_part_key(name)
        if part_key in parts_by_key:
            _fail("pptx_duplicate_zip_member")
        parts_by_key[part_key] = name
        if info.flag_bits & 0x1:
            _fail("pptx_encrypted_zip_member", _ENCRYPTED_USER_MESSAGE)
        if info.compress_type not in _SUPPORTED_ZIP_COMPRESSION:
            _fail("pptx_unsupported_zip_compression")

        with archive.open(info, "r") as member:
            while True:
                block = member.read(_READ_CHUNK_BYTES)
                if not block:
                    break
                total += len(block)
                if total > _MAX_INFLATED_BYTES:
                    _fail("pptx_inflated_size_limit", _INFLATED_LIMIT_USER_MESSAGE)
    return _PackageParts(by_key=parts_by_key, case_fallback_names={})


def _validate_member_name(name: str) -> None:
    if (
        not name
        or "\\" in name
        or name.startswith("/")
        or any(ord(character) < 0x20 or ord(character) == 0x7F for character in name)
    ):
        _fail("pptx_invalid_zip_member_name")
    _validate_percent_encoding(name, reason_code="pptx_invalid_zip_member_name")
    segments = name.split("/")
    for index, segment in enumerate(segments):
        if segment == "" and index == len(segments) - 1:
            continue
        if segment in {"", ".", ".."} or segment.endswith("."):
            _fail("pptx_invalid_zip_member_name")


def _decode_and_validate_zip_member_name(raw_name: bytes, *, flags: int) -> str:
    if not raw_name or any(byte < 0x20 or byte == 0x7F for byte in raw_name):
        _fail("pptx_invalid_zip_member_name")
    encoding = "utf-8" if flags & 0x800 else "cp437"
    try:
        name = raw_name.decode(encoding, errors="strict")
    except UnicodeDecodeError:
        _fail("pptx_invalid_zip_member_name")
    _validate_member_name(name)
    return name


def _validate_percent_encoding(value: str, *, reason_code: str) -> None:
    cursor = 0
    while cursor < len(value):
        character = value[cursor]
        if character != "%":
            if character.isspace():
                _fail(reason_code)
            cursor += 1
            continue
        if (
            cursor + 2 >= len(value)
            or value[cursor + 1] not in _ASCII_HEX_DIGITS
            or value[cursor + 2] not in _ASCII_HEX_DIGITS
        ):
            _fail(reason_code)
        encoded = int(value[cursor + 1 : cursor + 3], 16)
        if (
            encoded in _ASCII_UNRESERVED_BYTES
            or encoded < 0x20
            or encoded == 0x7F
            or encoded in {ord("/"), ord("\\")}
        ):
            _fail(reason_code)
        cursor += 3


def _parse_content_types(
    archive: zipfile.ZipFile,
    parts: _PackageParts,
) -> _ContentTypes:
    root = _read_xml(
        archive,
        parts,
        _CONTENT_TYPES_PART,
        missing_reason="pptx_missing_content_types",
    )
    namespace, local_name = _split_qname(root.tag)
    if local_name != "Types" or namespace not in _CONTENT_TYPES_NAMESPACES:
        _fail("pptx_invalid_content_types")

    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    for child in root:
        child_namespace, child_name = _split_qname(child.tag)
        if child_namespace != namespace:
            continue
        if child_name == "Default":
            extension = child.get("Extension", "").strip().casefold()
            content_type = child.get("ContentType", "").strip()
            if not extension or not content_type or extension in defaults:
                _fail("pptx_invalid_content_types")
            defaults[extension] = content_type
        elif child_name == "Override":
            raw_part_name = child.get("PartName", "").strip()
            content_type = child.get("ContentType", "").strip()
            if not raw_part_name or not content_type:
                _fail("pptx_invalid_content_types")
            part_name = _normalise_part_name(raw_part_name)
            part_key = _canonical_part_key(part_name)
            if part_key in overrides:
                _fail("pptx_invalid_content_types")
            overrides[part_key] = content_type
    return _ContentTypes(defaults=defaults, overrides=overrides)


def _parse_relationships(
    archive: zipfile.ZipFile,
    parts: _PackageParts,
    part_name: str,
    *,
    missing_reason: str,
) -> dict[str, _Relationship]:
    root = _read_xml(archive, parts, part_name, missing_reason=missing_reason)
    namespace, local_name = _split_qname(root.tag)
    if local_name != "Relationships" or namespace not in _PACKAGE_RELATIONSHIP_NAMESPACES:
        _fail("pptx_invalid_relationships")

    relationships: dict[str, _Relationship] = {}
    for child in root:
        child_namespace, child_name = _split_qname(child.tag)
        if child_namespace != namespace or child_name != "Relationship":
            continue
        relationship_id = child.get("Id", "").strip()
        relationship_type = child.get("Type", "").strip()
        target = child.get("Target", "").strip()
        target_mode = child.get("TargetMode")
        if target_mode is not None:
            target_mode = target_mode.strip()
        if (
            not relationship_id
            or not relationship_type
            or not target
            or relationship_id in relationships
            or target_mode not in {None, "Internal", "External"}
        ):
            _fail("pptx_invalid_relationships")
        relationships[relationship_id] = _Relationship(
            relationship_id=relationship_id,
            relationship_type=relationship_type,
            target=target,
            target_mode=target_mode,
        )
    return relationships


def _validate_presentation(
    archive: zipfile.ZipFile,
    parts: _PackageParts,
    content_types: _ContentTypes,
    presentation_part: str,
    *,
    is_strict: bool,
) -> bool:
    root = _read_xml(
        archive,
        parts,
        presentation_part,
        missing_reason="pptx_missing_presentation_part",
    )
    presentation_ns = _STRICT_PRESENTATION_NS if is_strict else _TRANSITIONAL_PRESENTATION_NS
    relationship_ns = _STRICT_RELATIONSHIP_NS if is_strict else _TRANSITIONAL_RELATIONSHIP_NS
    if root.tag != f"{{{presentation_ns}}}presentation":
        _fail("pptx_invalid_presentation_xml")

    slide_ids = root.findall(f"{{{presentation_ns}}}sldIdLst/{{{presentation_ns}}}sldId")
    if not slide_ids:
        return False

    relationships_part = _relationships_part_name(presentation_part)
    relationships = _parse_relationships(
        archive,
        parts,
        relationships_part,
        missing_reason="pptx_missing_presentation_relationships",
    )
    expected_slide_type = f"{relationship_ns}/slide"
    seen_relationship_ids: set[str] = set()
    seen_slide_ids: set[int] = set()
    seen_slide_parts: set[str] = set()
    parser_case_compatibility = parts.reference_requires_case_fallback(relationships_part)

    for slide_id in slide_ids:
        # ST_SlideId restricts XSD 1.0 unsignedInt: whitespace collapses, then only
        # ASCII decimal digits (without a sign) are in its lexical space.
        raw_slide_id = slide_id.get("id", "").strip(" \t\r\n")
        if not raw_slide_id or any(character not in "0123456789" for character in raw_slide_id):
            _fail("pptx_invalid_slide_id")
        numeric_slide_id = int(raw_slide_id, 10)
        if not 256 <= numeric_slide_id <= 2_147_483_647:
            _fail("pptx_invalid_slide_id")
        if numeric_slide_id in seen_slide_ids:
            _fail("pptx_duplicate_slide_id")
        seen_slide_ids.add(numeric_slide_id)

        relationship_id = slide_id.get(f"{{{relationship_ns}}}id", "").strip()
        if not relationship_id or relationship_id in seen_relationship_ids:
            _fail("pptx_invalid_slide_reference")
        seen_relationship_ids.add(relationship_id)

        relationship = relationships.get(relationship_id)
        if relationship is None or relationship.relationship_type != expected_slide_type:
            _fail("pptx_invalid_slide_reference")
        if relationship.target_mode == "External":
            _fail("pptx_external_slide_part")

        slide_part = _resolve_relationship_target(
            source_part=presentation_part,
            target=relationship.target,
        )
        slide_part_key = _canonical_part_key(slide_part)
        if slide_part_key in seen_slide_parts:
            _fail("pptx_duplicate_slide_part")
        seen_slide_parts.add(slide_part_key)
        if not parts.contains(slide_part):
            _fail("pptx_missing_slide_part")
        slide_case_fallback = parts.reference_requires_case_fallback(slide_part)
        parser_case_compatibility = slide_case_fallback or parser_case_compatibility
        slide_content_type = content_types.for_part(slide_part)
        if (
            slide_content_type is None
            or slide_content_type.casefold() != _SLIDE_CONTENT_TYPE.casefold()
        ):
            _fail("pptx_invalid_slide_content_type")
        parser_case_compatibility = (
            parser_case_compatibility or slide_content_type != _SLIDE_CONTENT_TYPE
        )
        slide_root = _read_xml(
            archive,
            parts,
            slide_part,
            missing_reason="pptx_missing_slide_part",
        )
        if slide_root.tag != f"{{{presentation_ns}}}sld":
            _fail("pptx_invalid_slide_xml")
    return parser_case_compatibility


def _validate_reachable_relationship_targets(
    archive: zipfile.ZipFile,
    parts: _PackageParts,
    content_types: _ContentTypes,
    root_relationships: dict[str, _Relationship],
    *,
    presentation_part: str,
) -> bool:
    """Preflight every reachable internal relationship without fetching external targets."""

    pending = [presentation_part]
    seen: set[str] = set()
    parser_case_compatibility = False

    def enqueue_internal_target(
        source_part: str | None,
        relationship: _Relationship,
    ) -> None:
        nonlocal parser_case_compatibility
        if relationship.target_mode == "External":
            return
        target_part = _resolve_relationship_target(
            source_part=source_part,
            target=relationship.target,
            allow_fragment=True,
        )
        if not parts.contains(target_part):
            _fail("pptx_missing_related_part")
        # Relationships parts are discovered by their source-part convention;
        # OPC does not permit another Relationship to target one directly.
        if _is_relationship_part_name(target_part):
            _fail("pptx_relationship_targets_relationship_part")
        if content_types.for_part(target_part) is None:
            _fail("pptx_missing_related_content_type")
        target_case_fallback = parts.reference_requires_case_fallback(target_part)
        parser_case_compatibility = target_case_fallback or parser_case_compatibility
        if _canonical_part_key(target_part) not in seen:
            pending.append(target_part)

    for relationship in root_relationships.values():
        enqueue_internal_target(None, relationship)

    while pending:
        source_part = pending.pop()
        source_key = _canonical_part_key(source_part)
        if source_key in seen:
            continue
        seen.add(source_key)
        relationships_part = _relationships_part_name(source_part)
        if not parts.contains(relationships_part):
            continue
        relationships_case_fallback = parts.reference_requires_case_fallback(relationships_part)
        parser_case_compatibility = relationships_case_fallback or parser_case_compatibility
        relationships = _parse_relationships(
            archive,
            parts,
            relationships_part,
            missing_reason="pptx_missing_related_part",
        )
        for relationship in relationships.values():
            enqueue_internal_target(source_part, relationship)
    return parser_case_compatibility


def _read_xml(
    archive: zipfile.ZipFile,
    parts: _PackageParts,
    part_name: str,
    *,
    missing_reason: str,
) -> ET.Element:
    if not parts.contains(part_name):
        _fail(missing_reason)
    with archive.open(parts.physical_name(part_name), "r") as member:
        payload = member.read(_MAX_CORE_XML_BYTES + 1)
    if len(payload) > _MAX_CORE_XML_BYTES:
        _fail("pptx_core_xml_size_limit", _PACKAGE_COMPLEXITY_USER_MESSAGE)
    try:
        _reject_unsafe_xml_declarations(payload)
    except _UnsafeXmlDeclarationError:
        _fail("pptx_unsafe_xml_declaration")
    except expat.ExpatError:
        _fail("pptx_invalid_xml")
    try:
        return ET.fromstring(payload)
    except (ET.ParseError, UnicodeError, ValueError):
        _fail("pptx_invalid_xml")


def _reject_unsafe_xml_declarations(payload: bytes) -> None:
    parser = expat.ParserCreate()

    def reject(*_args: object) -> None:
        raise _UnsafeXmlDeclarationError

    def reject_external(*_args: object) -> int:
        raise _UnsafeXmlDeclarationError

    parser.StartDoctypeDeclHandler = reject
    parser.EntityDeclHandler = reject
    parser.ExternalEntityRefHandler = reject_external
    parser.Parse(payload, True)


def _normalise_part_name(raw_part_name: str) -> str:
    if (
        not raw_part_name.startswith("/")
        or raw_part_name.startswith("//")
        or "\\" in raw_part_name
        or any(
            ord(character) < 0x20 or ord(character) == 0x7F
            for character in raw_part_name
        )
    ):
        _fail("pptx_invalid_part_name")
    try:
        parsed = urlsplit(raw_part_name)
    except (UnicodeError, ValueError):
        _fail("pptx_invalid_part_name")
    if (
        parsed.scheme
        or parsed.netloc
        or parsed.query
        or parsed.fragment
    ):
        _fail("pptx_invalid_part_name")
    part_name = parsed.path[1:]
    if not part_name or part_name.endswith("/"):
        _fail("pptx_invalid_part_name")
    _validate_percent_encoding(part_name, reason_code="pptx_invalid_part_name")
    if any(
        segment in {"", ".", ".."} or segment.endswith(".")
        for segment in part_name.split("/")
    ):
        _fail("pptx_invalid_part_name")
    return part_name


def _resolve_relationship_target(
    *,
    source_part: str | None,
    target: str,
    allow_fragment: bool = False,
) -> str:
    _validate_uri_reference_text(target)
    try:
        parsed = urlsplit(target)
    except (UnicodeError, ValueError):
        _fail("pptx_invalid_relationship_target")
    if (
        parsed.scheme
        or parsed.netloc
        or parsed.query
        or (parsed.fragment and not allow_fragment)
    ):
        _fail("pptx_invalid_relationship_target")
    raw_path = parsed.path
    if not raw_path and allow_fragment and parsed.fragment and source_part is not None:
        return _normalise_package_path(source_part, base_parts=[])
    if raw_path.startswith("/"):
        base_parts: list[str] = []
    elif source_part is None:
        base_parts = []
    else:
        base_parts = [part for part in posixpath.dirname(source_part).split("/") if part]
    return _normalise_package_path(raw_path, base_parts=base_parts)


def _normalise_package_path(raw_path: str, *, base_parts: list[str]) -> str:
    if (
        not raw_path
        or "\\" in raw_path
        or any(ord(character) < 0x20 or ord(character) == 0x7F for character in raw_path)
    ):
        _fail("pptx_invalid_package_path")
    _validate_percent_encoding(raw_path, reason_code="pptx_invalid_package_path")
    stack = list(base_parts)
    for segment in raw_path.split("/"):
        if segment in {"", "."}:
            continue
        if segment == "..":
            if not stack:
                _fail("pptx_package_path_escape")
            stack.pop()
        else:
            if segment.endswith("."):
                _fail("pptx_invalid_package_path")
            stack.append(segment)
    if not stack:
        _fail("pptx_invalid_package_path")
    return "/".join(stack)


def _validate_uri_reference_text(value: str) -> None:
    if any(ord(character) < 0x20 or ord(character) == 0x7F for character in value):
        _fail("pptx_invalid_package_path")


def _canonical_part_key(part_name: str) -> str:
    """Return the ASCII case-insensitive OPC equivalence key for a part name."""

    normalised = _normalise_package_path(part_name, base_parts=[])
    return normalised.translate(_ASCII_LOWER_TRANSLATION)


def _relationships_part_name(part_name: str) -> str:
    parent, filename = posixpath.split(part_name)
    return posixpath.join(parent, "_rels", f"{filename}.rels")


def _is_relationship_part_name(part_name: str) -> bool:
    """Return whether an OPC part name has the reserved Relationships-part shape."""

    normalised = _normalise_package_path(part_name, base_parts=[])
    segments = [segment.translate(_ASCII_LOWER_TRANSLATION) for segment in normalised.split("/")]
    if segments == ["_rels", ".rels"]:
        return True
    return (
        len(segments) >= 2
        and segments[-2] == "_rels"
        and segments[-1].endswith(".rels")
        and segments[-1] != ".rels"
    )


def _split_qname(tag: str) -> tuple[str, str]:
    if not tag.startswith("{") or "}" not in tag:
        return "", tag
    namespace, local_name = tag[1:].split("}", 1)
    return namespace, local_name


def _python_pptx_smoke_test(
    payload: bytes,
    *,
    allow_opc_case_compatibility: bool = False,
    case_fallback_names: dict[str, str] | None = None,
) -> ArtifactValidationReport:
    try:
        from pptx import Presentation
        from pptx.exc import InvalidXmlError, PackageNotFoundError
    except ImportError:
        _fail("pptx_parser_unavailable")

    def load(candidate: bytes) -> None:
        presentation = Presentation(io.BytesIO(candidate))
        for slide in presentation.slides:
            len(slide.shapes)

    try:
        load(payload)
    except MemoryError:
        raise
    except (NotImplementedError, AttributeError, TypeError):
        return ArtifactValidationReport(warnings=("pptx_parser_compatibility_warning",))
    except (PackageNotFoundError, KeyError, ValueError):
        if not allow_opc_case_compatibility:
            _fail("pptx_parser_structural_failure")
        compatibility_payload = _normalise_opc_case_for_parser(
            payload,
            case_fallback_names or {},
        )
        try:
            load(compatibility_payload)
        except MemoryError:
            raise
        except (NotImplementedError, AttributeError, TypeError):
            return ArtifactValidationReport(warnings=("pptx_parser_compatibility_warning",))
        except (
            PackageNotFoundError,
            KeyError,
            ValueError,
            InvalidXmlError,
            OSError,
            EOFError,
            zipfile.BadZipFile,
            SyntaxError,
        ):
            _fail("pptx_parser_structural_failure")
        except Exception:  # noqa: BLE001 - unknown parser incompatibilities remain soft
            return ArtifactValidationReport(warnings=("pptx_parser_compatibility_warning",))
        return ArtifactValidationReport(warnings=("pptx_opc_case_compatibility_warning",))
    except (
        InvalidXmlError,
        OSError,
        EOFError,
        zipfile.BadZipFile,
        SyntaxError,
    ):
        _fail("pptx_parser_structural_failure")
    except Exception:  # noqa: BLE001 - unknown parser incompatibilities are soft by contract
        return ArtifactValidationReport(warnings=("pptx_parser_compatibility_warning",))
    if allow_opc_case_compatibility:
        return ArtifactValidationReport(warnings=("pptx_opc_case_compatibility_warning",))
    return ArtifactValidationReport()


def _normalise_opc_case_for_parser(
    payload: bytes,
    case_fallback_names: dict[str, str],
) -> bytes:
    """Create an in-memory parser probe with only proven OPC case mismatches repaired."""

    output = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(payload), "r") as source,
        zipfile.ZipFile(output, "w") as destination,
    ):
        for info in source.infolist():
            member_payload = source.read(info)
            destination_name = case_fallback_names.get(info.filename, info.filename)
            if _canonical_part_key(info.filename) == _canonical_part_key(_CONTENT_TYPES_PART):
                root = ET.fromstring(member_payload)
                for child in root:
                    content_type = child.get("ContentType")
                    if content_type is None:
                        continue
                    folded_content_type = content_type.casefold()
                    if folded_content_type == _PRESENTATION_CONTENT_TYPE.casefold():
                        child.set("ContentType", _PRESENTATION_CONTENT_TYPE)
                    elif folded_content_type == _SLIDE_CONTENT_TYPE.casefold():
                        child.set("ContentType", _SLIDE_CONTENT_TYPE)
                member_payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            destination.writestr(
                destination_name,
                member_payload,
                compress_type=info.compress_type,
            )
    return output.getvalue()


def _fail(reason_code: str, user_message: str = _GENERIC_USER_MESSAGE) -> Never:
    raise ArtifactValidationError(reason_code, user_message)

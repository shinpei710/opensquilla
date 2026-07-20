from __future__ import annotations

import base64
import builtins
import copy
import io
import struct
import warnings
import xml.etree.ElementTree as ET
import zipfile
import zlib
from collections.abc import Callable
from typing import IO

import pytest
import structlog.testing
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.exc import PackageNotFoundError
from pptx.util import Inches

import opensquilla.artifact_validation as validation
from opensquilla.artifact_validation import (
    PPTX_MIME,
    ArtifactValidationError,
    ArtifactValidationReport,
    is_pptx_candidate,
    validate_artifact_for_delivery,
    validate_pptx_bytes,
)

_PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_TRANSITIONAL_REL_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
_STRICT_REL_NS = "http://purl.oclc.org/ooxml/officeDocument/relationships"
_TRANSITIONAL_PRESENTATION_NS = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)
_STRICT_PRESENTATION_NS = "http://purl.oclc.org/ooxml/presentationml/main"


class _UnseekableBytesIO(io.BytesIO):
    def seekable(self) -> bool:
        return False

    def seek(self, offset: int, whence: int = io.SEEK_SET) -> int:
        raise io.UnsupportedOperation


def _make_pptx(*, slides: int = 1, hyperlink: bool = False) -> bytes:
    presentation = Presentation()
    for index in range(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        run = text_box.text_frame.paragraphs[0].add_run()
        run.text = f"第 {index + 1} 页"
        if hyperlink:
            run.hyperlink.address = "https://example.invalid/presentation"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _entries(payload: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(io.BytesIO(payload), "r") as archive:
        return {info.filename: archive.read(info) for info in archive.infolist()}


def _pack(
    entries: dict[str, bytes],
    *,
    comment: bytes = b"",
    compression: int = zipfile.ZIP_DEFLATED,
) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)
        archive.comment = comment
    return output.getvalue()


def _pack_with_local_zip64_or_descriptors(
    entries: dict[str, bytes],
    *,
    streaming: bool,
    force_zip64: bool,
) -> bytes:
    output: io.BytesIO = _UnseekableBytesIO() if streaming else io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            if force_zip64:
                with archive.open(name, "w", force_zip64=True) as member:
                    member.write(data)
            else:
                archive.writestr(name, data)
    return output.getvalue()


def _case_variant_pptx() -> bytes:
    entries = _entries(_make_pptx())
    entries["ppt/slides/Slide1.xml"] = entries.pop("ppt/slides/slide1.xml")
    entries["ppt/slides/_rels/Slide1.xml.rels"] = entries.pop(
        "ppt/slides/_rels/slide1.xml.rels"
    )
    content_types = ET.fromstring(entries["[Content_Types].xml"])
    for child in content_types:
        if child.get("PartName") in {
            "/ppt/presentation.xml",
            "/ppt/slides/slide1.xml",
        }:
            child.set("ContentType", child.get("ContentType", "").upper())
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )
    return _pack(entries)


def _with_unindexed_local_bytes(payload: bytes, extra: bytes) -> bytes:
    eocd_offset = payload.rfind(b"PK\x05\x06")
    assert eocd_offset >= 0
    central_offset = struct.unpack_from("<L", payload, eocd_offset + 16)[0]
    modified = bytearray(payload[:central_offset] + extra + payload[central_offset:])
    struct.pack_into("<L", modified, eocd_offset + len(extra) + 16, central_offset + len(extra))
    return bytes(modified)


def _as_zip64(payload: bytes) -> bytes:
    eocd_offset = payload.rfind(b"PK\x05\x06")
    assert eocd_offset >= 0
    fields = struct.unpack_from("<4H2LH", payload, eocd_offset + 4)
    disk_number, central_disk, disk_entries, total_entries, size, offset, _ = fields
    zip64_record = b"PK\x06\x06" + struct.pack(
        "<Q2H2L4Q",
        44,
        45,
        45,
        disk_number,
        central_disk,
        disk_entries,
        total_entries,
        size,
        offset,
    )
    locator = b"PK\x06\x07" + struct.pack("<LQL", 0, eocd_offset, 1)
    legacy_eocd = bytearray(payload[eocd_offset:])
    struct.pack_into("<2H2L", legacy_eocd, 8, 0xFFFF, 0xFFFF, 0xFFFFFFFF, 0xFFFFFFFF)
    return payload[:eocd_offset] + zip64_record + locator + bytes(legacy_eocd)


def _rewrite_xml(
    payload: bytes,
    part_name: str,
    mutate: Callable[[ET.Element], None],
) -> bytes:
    entries = _entries(payload)
    root = ET.fromstring(entries[part_name])
    mutate(root)
    entries[part_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    return _pack(entries)


def _remove_part(payload: bytes, part_name: str) -> bytes:
    entries = _entries(payload)
    entries.pop(part_name)
    return _pack(entries)


def _validate(payload: bytes, **overrides: str | None) -> ArtifactValidationReport:
    arguments: dict[str, str | None] = {
        "source_name": "generated.pptx",
        "name": "generated.pptx",
        "mime": PPTX_MIME,
        "source": "test",
    }
    arguments.update(overrides)
    return validate_artifact_for_delivery(
        payload,
        source_name=arguments["source_name"],
        name=str(arguments["name"]),
        mime=str(arguments["mime"]),
        source=str(arguments["source"]),
    )


@pytest.mark.parametrize(
    ("source_name", "name", "mime"),
    [
        ("DECK.PPTX", "deck.bin", "application/octet-stream"),
        ("deck.bin", "RESULT.PpTx", "application/octet-stream"),
        (None, "deck.bin", f" {PPTX_MIME.upper()}; charset=binary"),
    ],
)
def test_is_pptx_candidate_uses_any_delivery_signal(
    source_name: str | None,
    name: str,
    mime: str,
) -> None:
    assert is_pptx_candidate(source_name=source_name, name=name, mime=mime)


def test_non_pptx_is_an_opaque_noop() -> None:
    assert validate_artifact_for_delivery(
        b"not an office document",
        source_name="report.bin",
        name="report.bin",
        mime="application/octet-stream",
        source="test",
    ) == ArtifactValidationReport()


@pytest.mark.parametrize("slide_count", [0, 1, 3])
def test_valid_transitional_pptx_is_deep_loaded(slide_count: int) -> None:
    assert _validate(_make_pptx(slides=slide_count)) == ArtifactValidationReport()


def test_external_hyperlinks_and_unknown_parts_are_allowed() -> None:
    entries = _entries(_make_pptx(hyperlink=True))
    entries["custom/unknown.bin"] = b"extension data"
    entries["custom/兼容扩展.bin"] = b"tolerated non-canonical unknown entry"

    assert _validate(_pack(entries)) == ArtifactValidationReport()


def test_internal_relationship_fragment_is_allowed_without_network_access() -> None:
    entries = _entries(_make_pptx())
    relationships = ET.fromstring(entries["ppt/slides/_rels/slide1.xml.rels"])
    ET.SubElement(
        relationships,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rIdInternalLink",
            "Type": f"{_TRANSITIONAL_REL_NS}/hyperlink",
            "Target": "#section-1",
            "TargetMode": "Internal",
        },
    )
    entries["ppt/slides/_rels/slide1.xml.rels"] = ET.tostring(
        relationships,
        encoding="utf-8",
        xml_declaration=True,
    )

    assert _validate(_pack(entries)) == ArtifactValidationReport()


def test_complex_picture_chart_and_notes_pptx_is_accepted() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUB"
        "AScY42YAAAAASUVORK5CYII="
    )
    slide.shapes.add_picture(io.BytesIO(png), Inches(0.5), Inches(0.5), Inches(1), Inches(1))
    chart_data = ChartData()
    chart_data.categories = ["一月", "二月"]
    chart_data.add_series("收入", (3, 5))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(2),
        Inches(1),
        Inches(5),
        Inches(3),
        chart_data,
    )
    slide.notes_slide.notes_text_frame.text = "仅供讲者查看的备注"
    output = io.BytesIO()
    presentation.save(output)

    assert validate_pptx_bytes(output.getvalue()) == ArtifactValidationReport()


def test_zip_comment_is_allowed_when_eocd_reaches_exact_end() -> None:
    assert _validate(_pack(_entries(_make_pptx()), comment=b"valid comment")) == (
        ArtifactValidationReport()
    )


def test_valid_zip64_directory_is_supported() -> None:
    assert _validate(_as_zip64(_make_pptx())) == ArtifactValidationReport()


@pytest.mark.parametrize(
    ("streaming", "force_zip64"),
    [
        (True, False),
        (False, True),
        (True, True),
    ],
)
def test_valid_data_descriptors_and_entry_level_zip64_are_supported(
    streaming: bool,
    force_zip64: bool,
) -> None:
    payload = _pack_with_local_zip64_or_descriptors(
        _entries(_make_pptx()),
        streaming=streaming,
        force_zip64=force_zip64,
    )

    assert _validate(payload) == ArtifactValidationReport()


def test_unsigned_data_descriptor_crc_may_equal_optional_signature() -> None:
    signature_as_crc = struct.unpack("<L", b"PK\x07\x08")[0]
    entry = validation._CentralDirectoryEntry(
        decoded_name="probe.bin",
        raw_name=b"probe.bin",
        version_needed=20,
        flags=0x08,
        compression=zipfile.ZIP_DEFLATED,
        modified_time=0,
        modified_date=0,
        crc=signature_as_crc,
        compressed_size=4,
        uncompressed_size=4,
        local_header_offset=0,
        uses_zip64_sizes=False,
    )
    descriptor = struct.pack("<LLL", signature_as_crc, 4, 4)

    validation._validate_data_descriptor(descriptor, entry, uses_zip64=False)


def test_zip64_sentinel_without_directory_is_rejected() -> None:
    payload = bytearray(_make_pptx())
    eocd_offset = payload.rfind(b"PK\x05\x06")
    struct.pack_into("<2H", payload, eocd_offset + 8, 0xFFFF, 0xFFFF)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(payload))
    assert raised.value.reason_code == "pptx_invalid_zip64_directory"


def test_central_directory_must_end_immediately_before_eocd() -> None:
    payload = _make_pptx()
    eocd_offset = payload.rfind(b"PK\x05\x06")

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload[:eocd_offset] + b"unexpected-gap" + payload[eocd_offset:])
    assert raised.value.reason_code == "pptx_invalid_central_directory"


def test_member_count_is_bounded_before_zipfile_parsing(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(validation, "_MAX_ZIP_MEMBERS", 2)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_make_pptx())
    assert raised.value.reason_code == "pptx_member_count_limit"


@pytest.mark.parametrize(
    ("payload", "reason_code"),
    [
        (b"not a zip", "pptx_not_zip"),
        (b"prefix" + _make_pptx(), "pptx_not_zip"),
        (_make_pptx()[:-10], "pptx_invalid_zip_envelope"),
        (_make_pptx() + b"trailing bytes", "pptx_invalid_zip_envelope"),
    ],
    ids=("plain-non-zip", "prefixed-zip", "truncated-zip", "trailing-bytes"),
)
def test_non_pure_or_truncated_zip_is_rejected(payload: bytes, reason_code: str) -> None:
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload)

    assert raised.value.reason_code == reason_code
    assert "generated.pptx" not in raised.value.user_message


def test_ole_or_encrypted_container_has_actionable_safe_error() -> None:
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"encrypted")

    assert raised.value.reason_code == "pptx_encrypted_or_legacy_container"
    assert "unencrypted .pptx" in raised.value.user_message


def test_duplicate_zip_member_is_rejected() -> None:
    payload = _make_pptx()
    output = io.BytesIO()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", UserWarning)
        with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
            for name, data in _entries(payload).items():
                archive.writestr(name, data)
            archive.writestr("[Content_Types].xml", b"duplicate")

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(output.getvalue())
    assert raised.value.reason_code == "pptx_duplicate_zip_member"


@pytest.mark.parametrize("compression", [zipfile.ZIP_BZIP2, zipfile.ZIP_LZMA])
def test_only_stored_and_deflated_zip_members_are_accepted(compression: int) -> None:
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(_entries(_make_pptx()), compression=compression))
    assert raised.value.reason_code == "pptx_unsupported_zip_compression"


@pytest.mark.parametrize(
    "invalid_name",
    [
        "custom/bad%.bin",
        "custom/%41.bin",
        "custom/%2f.bin",
        "custom/%5C.bin",
        "custom/%00.bin",
        "custom/%7f.bin",
        "custom/a%41.bin",
        "custom/ab%41.bin",
    ],
)
def test_zip_member_names_follow_pack_uri_percent_encoding_rules(
    invalid_name: str,
) -> None:
    entries = _entries(_make_pptx())
    entries[invalid_name] = b"invalid OPC part name"

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_zip_member_name"


@pytest.mark.parametrize("raw_control", [b"\x00", b"\x01", b"\x7f"])
def test_raw_central_directory_member_name_controls_are_rejected(
    raw_control: bytes,
) -> None:
    entries = _entries(_make_pptx())
    marker = b"custom/probe.bin"
    entries[marker.decode()] = b"probe"
    payload = _pack(entries)
    assert payload.count(marker) == 2
    replacement = marker[:7] + raw_control + marker[8:]

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload.replace(marker, replacement))
    assert raised.value.reason_code == "pptx_invalid_zip_member_name"


def test_percent_escape_hex_case_is_part_name_equivalent() -> None:
    entries = _entries(_make_pptx())
    entries["custom/my%2A.bin"] = b"first"
    entries["custom/my%2a.bin"] = b"second"

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_duplicate_zip_member"


def test_local_header_filename_must_match_central_directory() -> None:
    entries = _entries(_make_pptx())
    marker = b"custom/probe.bin"
    entries[marker.decode()] = b"probe"
    payload = bytearray(_pack(entries))
    local_name_offset = payload.find(marker)
    central_name_offset = payload.rfind(marker)
    assert 0 <= local_name_offset < central_name_offset
    payload[local_name_offset] = ord("C")

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(payload))
    assert raised.value.reason_code == "pptx_invalid_local_file_header"


def test_local_header_compression_must_match_central_directory() -> None:
    payload = bytearray(_make_pptx())
    assert payload.startswith(b"PK\x03\x04")
    central_method = struct.unpack_from("<H", payload, payload.find(b"PK\x01\x02") + 10)[0]
    replacement = zipfile.ZIP_STORED if central_method == zipfile.ZIP_DEFLATED else 99
    struct.pack_into("<H", payload, 8, replacement)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(payload))
    assert raised.value.reason_code == "pptx_invalid_local_file_header"


def test_local_header_crc_must_match_central_directory() -> None:
    payload = bytearray(_make_pptx())
    local_crc = struct.unpack_from("<L", payload, 14)[0]
    struct.pack_into("<L", payload, 14, local_crc ^ 0xFFFFFFFF)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(payload))
    assert raised.value.reason_code == "pptx_invalid_local_file_header"


def test_unindexed_local_record_or_gap_is_rejected() -> None:
    unindexed_record = b"PK\x03\x04" + (b"\x00" * 26)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_with_unindexed_local_bytes(_make_pptx(), unindexed_record))
    assert raised.value.reason_code == "pptx_invalid_local_file_layout"


def test_overlapping_local_record_offsets_are_rejected() -> None:
    payload = bytearray(_make_pptx())
    eocd_offset = payload.rfind(b"PK\x05\x06")
    central_offset = struct.unpack_from("<L", payload, eocd_offset + 16)[0]
    first_local_offset = struct.unpack_from("<L", payload, central_offset + 42)[0]
    name_length, extra_length, comment_length = struct.unpack_from(
        "<3H", payload, central_offset + 28
    )
    second_central_entry = central_offset + 46 + name_length + extra_length + comment_length
    assert payload[second_central_entry : second_central_entry + 4] == b"PK\x01\x02"
    struct.pack_into("<L", payload, second_central_entry + 42, first_local_offset)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(payload))
    assert raised.value.reason_code == "pptx_invalid_local_file_layout"


def test_crc_failure_is_rejected() -> None:
    entries = _entries(_make_pptx())
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)
        archive.writestr("probe.bin", b"crc payload", compress_type=zipfile.ZIP_STORED)
    corrupted = bytearray(output.getvalue())
    marker = b"probe.bin"
    header_offset = corrupted.find(b"PK\x03\x04", corrupted.find(marker) - 30)
    filename_length, extra_length = struct.unpack_from("<HH", corrupted, header_offset + 26)
    data_offset = header_offset + 30 + filename_length + extra_length
    corrupted[data_offset] ^= 0x01

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(bytes(corrupted))
    assert raised.value.reason_code == "pptx_corrupt_zip_member"


def test_raw_zlib_failures_are_classified_as_corrupt_members(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    payload = _make_pptx()

    def fail_open(*_args: object, **_kwargs: object) -> object:
        raise zlib.error("synthetic damaged deflate stream")

    monkeypatch.setattr(zipfile.ZipFile, "open", fail_open)
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload)
    assert raised.value.reason_code == "pptx_corrupt_zip_member"


def test_actual_inflated_bytes_are_bounded(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(validation, "_MAX_INFLATED_BYTES", 64)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_make_pptx())
    assert raised.value.reason_code == "pptx_inflated_size_limit"


def test_core_xml_member_size_is_bounded(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(validation, "_MAX_CORE_XML_BYTES", 64)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_make_pptx())
    assert raised.value.reason_code == "pptx_core_xml_size_limit"


@pytest.mark.parametrize(
    ("part_name", "reason_code"),
    [
        ("[Content_Types].xml", "pptx_missing_content_types"),
        ("_rels/.rels", "pptx_missing_root_relationships"),
        ("ppt/presentation.xml", "pptx_missing_presentation_part"),
    ],
)
def test_required_opc_parts_must_exist(part_name: str, reason_code: str) -> None:
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_remove_part(_make_pptx(), part_name))
    assert raised.value.reason_code == reason_code


def test_core_xml_must_parse() -> None:
    entries = _entries(_make_pptx())
    entries["[Content_Types].xml"] = b"<Types"

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_xml"


def test_doctype_text_inside_cdata_is_not_treated_as_a_declaration() -> None:
    entries = _entries(_make_pptx())
    entries["ppt/slides/slide1.xml"] = entries["ppt/slides/slide1.xml"].replace(
        "第 1 页".encode(),
        b"<![CDATA[<!DOCTYPE html>]]>",
    )
    payload = _pack(entries)

    assert _validate(payload) == ArtifactValidationReport()
    reloaded = Presentation(io.BytesIO(payload))
    assert reloaded.slides[0].shapes[0].text == "<!DOCTYPE html>"


def test_real_doctype_declaration_is_rejected() -> None:
    entries = _entries(_make_pptx())
    content_types = entries["[Content_Types].xml"]
    declaration_end = content_types.find(b"?>") + 2
    entries["[Content_Types].xml"] = (
        content_types[:declaration_end]
        + b"<!DOCTYPE Types>"
        + content_types[declaration_end:]
    )

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_unsafe_xml_declaration"


def test_root_must_have_one_office_document_relationship() -> None:
    def duplicate_office_relationship(root: ET.Element) -> None:
        office_relationship = next(
            child for child in root if child.get("Type", "").endswith("/officeDocument")
        )
        clone = copy.deepcopy(office_relationship)
        clone.set("Id", "duplicateOfficeDocument")
        root.append(clone)

    payload = _rewrite_xml(_make_pptx(), "_rels/.rels", duplicate_office_relationship)
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload)
    assert raised.value.reason_code == "pptx_invalid_office_document_relationship"


def test_explicit_internal_target_mode_is_accepted() -> None:
    def mark_internal(root: ET.Element) -> None:
        office_relationship = next(
            child for child in root if child.get("Type", "").endswith("/officeDocument")
        )
        office_relationship.set("TargetMode", "Internal")

    assert _validate(_rewrite_xml(_make_pptx(), "_rels/.rels", mark_internal)) == (
        ArtifactValidationReport()
    )


@pytest.mark.parametrize(
    "presentation_part",
    [
        "custom/presentation.xml",
        "custom/presentation.v1.xml",
        "custom/my%20deck.xml",
    ],
)
@pytest.mark.parametrize("strict", [False, True])
def test_main_presentation_part_can_use_nonstandard_opc_path(
    presentation_part: str,
    strict: bool,
) -> None:
    entries = _entries(_make_pptx())
    main_xml = entries.pop("ppt/presentation.xml")
    main_rels = ET.fromstring(entries.pop("ppt/_rels/presentation.xml.rels"))
    for relationship in main_rels:
        if relationship.get("TargetMode") != "External":
            relationship.set("Target", f"../ppt/{relationship.get('Target')}")
    entries[presentation_part] = main_xml
    relationships_part = (
        f"custom/_rels/{presentation_part.rsplit('/', 1)[-1]}.rels"
    )
    entries[relationships_part] = ET.tostring(
        main_rels,
        encoding="utf-8",
        xml_declaration=True,
    )

    root_rels = ET.fromstring(entries["_rels/.rels"])
    office_relationship = next(
        child for child in root_rels if child.get("Type", "").endswith("/officeDocument")
    )
    office_relationship.set("Target", presentation_part)
    entries["_rels/.rels"] = ET.tostring(root_rels, encoding="utf-8", xml_declaration=True)

    content_types = ET.fromstring(entries["[Content_Types].xml"])
    main_override = next(
        child
        for child in content_types
        if child.get("PartName") == "/ppt/presentation.xml"
    )
    main_override.set("PartName", f"/{presentation_part}")
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )

    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    report = _validate(_pack(entries))
    expected_warnings = ("pptx_strict_ooxml_smoke_skipped",) if strict else ()
    assert report.warnings == expected_warnings


@pytest.mark.parametrize("strict", [False, True])
@pytest.mark.parametrize(
    "invalid_part_name",
    [
        "//ppt/slides/slide1.xml",
        "/ppt//slides/slide1.xml",
        "/ppt/slides/./slide1.xml",
        "/ppt/slides/../slides/slide1.xml",
        "/ppt/slides/slide1.xml/",
        "/ppt/slides/slide1.xml.",
        "/ppt\\slides\\slide1.xml",
        "/ppt/slides/slide%ZZ.xml",
    ],
)
def test_content_type_override_part_name_is_not_path_normalised(
    strict: bool,
    invalid_part_name: str,
) -> None:
    entries = _entries(_make_pptx())
    content_types = ET.fromstring(entries["[Content_Types].xml"])
    slide_override = next(
        child
        for child in content_types
        if child.get("PartName") == "/ppt/slides/slide1.xml"
    )
    slide_override.set("PartName", invalid_part_name)
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_part_name"


@pytest.mark.parametrize("strict", [False, True])
def test_relationship_target_segment_cannot_end_with_dot(strict: bool) -> None:
    def use_trailing_dot_target(root: ET.Element) -> None:
        slide_relationship = next(
            child for child in root if child.get("Type", "").endswith("/slide")
        )
        slide_relationship.set("Target", "slides/missing.xml.")

    entries = _entries(
        _rewrite_xml(
            _make_pptx(),
            "ppt/_rels/presentation.xml.rels",
            use_trailing_dot_target,
        )
    )
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_package_path"


@pytest.mark.parametrize("member_name", ["custom/data.bin.", "custom/folder./"])
def test_zip_member_segment_cannot_end_with_dot(member_name: str) -> None:
    entries = _entries(_make_pptx())
    entries[member_name] = b"ambiguous on Windows"

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_zip_member_name"


@pytest.mark.parametrize("strict", [False, True])
def test_synchronised_trailing_dot_slide_part_is_rejected(strict: bool) -> None:
    entries = _entries(_make_pptx())
    entries["ppt/slides/slide1.xml."] = entries.pop("ppt/slides/slide1.xml")
    entries["ppt/slides/_rels/slide1.xml..rels"] = entries.pop(
        "ppt/slides/_rels/slide1.xml.rels"
    )

    presentation_relationships = ET.fromstring(
        entries["ppt/_rels/presentation.xml.rels"]
    )
    slide_relationship = next(
        child
        for child in presentation_relationships
        if child.get("Type", "").endswith("/slide")
    )
    slide_relationship.set("Target", "slides/slide1.xml.")
    entries["ppt/_rels/presentation.xml.rels"] = ET.tostring(
        presentation_relationships,
        encoding="utf-8",
        xml_declaration=True,
    )

    content_types = ET.fromstring(entries["[Content_Types].xml"])
    slide_override = next(
        child
        for child in content_types
        if child.get("PartName") == "/ppt/slides/slide1.xml"
    )
    slide_override.set("PartName", "/ppt/slides/slide1.xml.")
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_invalid_zip_member_name"


def test_opc_part_and_content_type_case_variants_still_run_parser_smoke(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import pptx

    parser_calls = 0
    real_presentation = pptx.Presentation

    def record_parser_call(stream: str | IO[bytes] | None) -> object:
        nonlocal parser_calls
        parser_calls += 1
        return real_presentation(stream)

    monkeypatch.setattr(pptx, "Presentation", record_parser_call)

    assert _validate(_case_variant_pptx()).warnings == (
        "pptx_opc_case_compatibility_warning",
    )
    assert parser_calls >= 1


@pytest.mark.parametrize(
    "parser_error",
    [
        PackageNotFoundError("unrelated package failure"),
        ValueError("unrelated structural failure"),
        KeyError("unrelated missing part"),
    ],
)
def test_opc_case_compatibility_does_not_hide_unrelated_parser_failures(
    monkeypatch: pytest.MonkeyPatch,
    parser_error: Exception,
) -> None:
    import pptx

    def fail_to_open(_stream: object) -> object:
        raise parser_error

    monkeypatch.setattr(pptx, "Presentation", fail_to_open)
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_case_variant_pptx())
    assert raised.value.reason_code == "pptx_parser_structural_failure"


def test_ascii_case_equivalent_zip_members_are_duplicates() -> None:
    entries = _entries(_make_pptx())
    entries["PPT/PRESENTATION.XML"] = entries["ppt/presentation.xml"]

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_duplicate_zip_member"


def test_strict_ooxml_slide_graph_is_structurally_checked_and_soft_warned() -> None:
    entries = _entries(_make_pptx(slides=1))
    for part_name, payload in list(entries.items()):
        if part_name.endswith((".xml", ".rels")):
            entries[part_name] = payload.replace(
                _TRANSITIONAL_PRESENTATION_NS.encode(),
                _STRICT_PRESENTATION_NS.encode(),
            ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    report = _validate(_pack(entries))
    assert report.warnings == ("pptx_strict_ooxml_smoke_skipped",)


def test_broken_slide_relationship_is_rejected() -> None:
    def break_slide_target(root: ET.Element) -> None:
        slide_relationship = next(
            child for child in root if child.get("Type", "").endswith("/slide")
        )
        slide_relationship.set("Target", "slides/missing.xml")

    payload = _rewrite_xml(
        _make_pptx(),
        "ppt/_rels/presentation.xml.rels",
        break_slide_target,
    )
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload)
    assert raised.value.reason_code == "pptx_missing_slide_part"


@pytest.mark.parametrize(
    "slide_id",
    [None, "", "one", "+256", "-1", "255", "2147483648"],
)
def test_slide_id_must_be_a_decimal_in_the_ooxml_range(slide_id: str | None) -> None:
    def replace_slide_id(root: ET.Element) -> None:
        node = root.find(
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldIdLst/"
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldId"
        )
        assert node is not None
        if slide_id is None:
            node.attrib.pop("id", None)
        else:
            node.set("id", slide_id)

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_rewrite_xml(_make_pptx(), "ppt/presentation.xml", replace_slide_id))
    assert raised.value.reason_code == "pptx_invalid_slide_id"


def test_slide_id_uses_xsd_whitespace_collapse() -> None:
    def pad_slide_id(root: ET.Element) -> None:
        node = root.find(
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldIdLst/"
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldId"
        )
        assert node is not None
        node.set("id", " \t256\n ")

    assert _validate(
        _rewrite_xml(_make_pptx(), "ppt/presentation.xml", pad_slide_id)
    ) == ArtifactValidationReport()


def test_slide_ids_are_unique_as_numeric_values() -> None:
    def duplicate_slide_id(root: ET.Element) -> None:
        nodes = root.findall(
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldIdLst/"
            f"{{{_TRANSITIONAL_PRESENTATION_NS}}}sldId"
        )
        assert len(nodes) == 2
        nodes[1].set("id", f"0{nodes[0].get('id')}")

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(
            _rewrite_xml(_make_pptx(slides=2), "ppt/presentation.xml", duplicate_slide_id)
        )
    assert raised.value.reason_code == "pptx_duplicate_slide_id"


def test_slide_list_targets_must_resolve_to_distinct_parts() -> None:
    def duplicate_slide_target(root: ET.Element) -> None:
        slide_relationships = [
            child for child in root if child.get("Type", "").endswith("/slide")
        ]
        assert len(slide_relationships) == 2
        slide_relationships[1].set("Target", slide_relationships[0].get("Target", ""))

    payload = _rewrite_xml(
        _make_pptx(slides=2),
        "ppt/_rels/presentation.xml.rels",
        duplicate_slide_target,
    )
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(payload)
    assert raised.value.reason_code == "pptx_duplicate_slide_part"


@pytest.mark.parametrize("strict", [False, True])
def test_every_reachable_internal_part_requires_a_content_type(strict: bool) -> None:
    entries = _entries(_make_pptx())
    slide_relationships = ET.fromstring(entries["ppt/slides/_rels/slide1.xml.rels"])
    ET.SubElement(
        slide_relationships,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rIdMissingContentType",
            "Type": f"{_TRANSITIONAL_REL_NS}/customData",
            "Target": "../custom/data.missingct",
        },
    )
    entries["ppt/slides/_rels/slide1.xml.rels"] = ET.tostring(
        slide_relationships,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["ppt/custom/data.missingct"] = b"custom data"
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_missing_related_content_type"


def test_root_relationship_targets_also_require_a_content_type() -> None:
    entries = _entries(_make_pptx())
    root_relationships = ET.fromstring(entries["_rels/.rels"])
    ET.SubElement(
        root_relationships,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rIdMissingContentType",
            "Type": f"{_TRANSITIONAL_REL_NS}/customData",
            "Target": "custom/data.missingct",
        },
    )
    entries["_rels/.rels"] = ET.tostring(
        root_relationships,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["custom/data.missingct"] = b"custom data"

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_missing_related_content_type"


@pytest.mark.parametrize("strict", [False, True])
def test_relationship_parts_do_not_require_content_types_coverage(strict: bool) -> None:
    entries = _entries(_make_pptx())
    content_types = ET.fromstring(entries["[Content_Types].xml"])
    rels_default = next(
        child
        for child in content_types
        if child.get("Extension", "").casefold() == "rels"
    )
    content_types.remove(rels_default)
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    report = _validate(_pack(entries))
    expected_warnings = ("pptx_strict_ooxml_smoke_skipped",) if strict else ()
    assert report.warnings == expected_warnings


@pytest.mark.parametrize("strict", [False, True])
def test_ordinary_rels_suffix_part_still_requires_a_content_type(strict: bool) -> None:
    entries = _entries(_make_pptx())
    content_types = ET.fromstring(entries["[Content_Types].xml"])
    rels_default = next(
        child
        for child in content_types
        if child.get("Extension", "").casefold() == "rels"
    )
    content_types.remove(rels_default)
    entries["[Content_Types].xml"] = ET.tostring(
        content_types,
        encoding="utf-8",
        xml_declaration=True,
    )

    slide_relationships = ET.fromstring(entries["ppt/slides/_rels/slide1.xml.rels"])
    ET.SubElement(
        slide_relationships,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rIdOrdinaryRelsPart",
            "Type": f"{_TRANSITIONAL_REL_NS}/customData",
            "Target": "../custom/data.rels",
        },
    )
    entries["ppt/slides/_rels/slide1.xml.rels"] = ET.tostring(
        slide_relationships,
        encoding="utf-8",
        xml_declaration=True,
    )
    entries["ppt/custom/data.rels"] = b"ordinary part, not a Relationships part"

    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_missing_related_content_type"


@pytest.mark.parametrize("strict", [False, True])
def test_relationship_cannot_target_a_relationships_part(strict: bool) -> None:
    entries = _entries(_make_pptx())
    slide_relationships = ET.fromstring(entries["ppt/slides/_rels/slide1.xml.rels"])
    ET.SubElement(
        slide_relationships,
        f"{{{_PACKAGE_REL_NS}}}Relationship",
        {
            "Id": "rIdRelationshipsPartTarget",
            "Type": f"{_TRANSITIONAL_REL_NS}/customData",
            "Target": "_rels/slide1.xml.rels",
        },
    )
    entries["ppt/slides/_rels/slide1.xml.rels"] = ET.tostring(
        slide_relationships,
        encoding="utf-8",
        xml_declaration=True,
    )
    if strict:
        for part_name, part_payload in list(entries.items()):
            if part_name.endswith((".xml", ".rels")):
                entries[part_name] = part_payload.replace(
                    _TRANSITIONAL_PRESENTATION_NS.encode(),
                    _STRICT_PRESENTATION_NS.encode(),
                ).replace(_TRANSITIONAL_REL_NS.encode(), _STRICT_REL_NS.encode())

    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_pack(entries))
    assert raised.value.reason_code == "pptx_relationship_targets_relationship_part"


def test_slide_content_type_and_root_are_checked() -> None:
    def change_content_type(root: ET.Element) -> None:
        override = next(
            child for child in root if child.get("PartName") == "/ppt/slides/slide1.xml"
        )
        override.set("ContentType", "application/xml")

    invalid_content_type = _rewrite_xml(
        _make_pptx(),
        "[Content_Types].xml",
        change_content_type,
    )
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(invalid_content_type)
    assert raised.value.reason_code == "pptx_invalid_slide_content_type"

    def change_root(root: ET.Element) -> None:
        root.tag = f"{{{_TRANSITIONAL_PRESENTATION_NS}}}not-a-slide"

    invalid_root = _rewrite_xml(_make_pptx(), "ppt/slides/slide1.xml", change_root)
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(invalid_root)
    assert raised.value.reason_code == "pptx_invalid_slide_xml"


@pytest.mark.parametrize(
    ("parser_error", "expected_warning"),
    [
        (NotImplementedError("unsupported feature"), True),
        (AttributeError("unknown parser shape"), True),
        (TypeError("unknown parser type"), True),
        (RuntimeError("unknown parser compatibility"), True),
        (ValueError("structural parser failure"), False),
        (KeyError("missing parser part"), False),
    ],
)
def test_python_pptx_parser_error_classification(
    monkeypatch: pytest.MonkeyPatch,
    parser_error: Exception,
    expected_warning: bool,
) -> None:
    import pptx

    def fail_to_open(_stream: object) -> object:
        raise parser_error

    monkeypatch.setattr(pptx, "Presentation", fail_to_open)
    if expected_warning:
        assert _validate(_make_pptx()).warnings == ("pptx_parser_compatibility_warning",)
    else:
        with pytest.raises(ArtifactValidationError) as raised:
            _validate(_make_pptx())
        assert raised.value.reason_code == "pptx_parser_structural_failure"


def test_python_pptx_import_failure_blocks_delivery(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    real_import = builtins.__import__

    def fail_pptx_import(
        name: str,
        globals: dict[str, object] | None = None,
        locals: dict[str, object] | None = None,
        fromlist: tuple[str, ...] = (),
        level: int = 0,
    ) -> object:
        if name == "pptx" or name.startswith("pptx."):
            raise ImportError("synthetic missing parser")
        return real_import(name, globals, locals, fromlist, level)

    monkeypatch.setattr(builtins, "__import__", fail_pptx_import)
    with pytest.raises(ArtifactValidationError) as raised:
        _validate(_make_pptx())
    assert raised.value.reason_code == "pptx_parser_unavailable"


def test_validation_logs_only_safe_bounded_metadata() -> None:
    with structlog.testing.capture_logs() as captured:
        with pytest.raises(ArtifactValidationError):
            validate_artifact_for_delivery(
                b"invalid",
                source_name="/secret/workspace/customer-deck.pptx",
                name="private-title.pptx",
                mime=PPTX_MIME,
                source="publish_artifact",
            )

    assert len(captured) == 1
    event = captured[0]
    assert event["source"] == "publish_artifact"
    assert event["outcome"] == "blocked"
    assert event["reason_code"] == "pptx_not_zip"
    assert event["size"] == 7
    assert "secret" not in repr(event)
    assert "private-title" not in repr(event)

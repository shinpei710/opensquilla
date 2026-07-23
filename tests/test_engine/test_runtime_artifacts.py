from __future__ import annotations

import asyncio
import contextlib
import hashlib
import json
import os
import threading
from collections.abc import AsyncIterator
from io import BytesIO
from types import SimpleNamespace
from typing import Any

import pytest
from pptx import Presentation

from opensquilla.artifacts import ArtifactStore, artifact_payload
from opensquilla.engine.artifact_delivery import (
    artifact_delivery_publish_target_key,
    auto_publish_omitted_workspace_artifacts,
)
from opensquilla.engine.runtime import TurnRunner
from opensquilla.engine.types import ArtifactEvent, DoneEvent, TextDeltaEvent, ToolUseStartEvent
from opensquilla.gateway.config import AttachmentsConfig, GatewayConfig, SquillaRouterConfig
from opensquilla.provider import DoneEvent as ProviderDone
from opensquilla.provider import Message, ModelInfo
from opensquilla.provider import TextDeltaEvent as ProviderText
from opensquilla.provider import ToolUseEndEvent as ProviderToolUseEnd
from opensquilla.provider import ToolUseStartEvent as ProviderToolUseStart
from opensquilla.session.manager import SessionManager
from opensquilla.session.storage import SessionStorage
from opensquilla.tools.builtin import filesystem
from opensquilla.tools.builtin import patch as patch_tools
from opensquilla.tools.registry import ToolRegistry, ToolSpec
from opensquilla.tools.types import (
    CallerKind,
    RetryableToolInputError,
    ToolContext,
    ToolError,
    current_tool_context,
)


class _ArtifactProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="tool-1", tool_name="make_file")
            yield ProviderToolUseEnd(
                tool_use_id="tool-1",
                tool_name="make_file",
                arguments={},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="done")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _BlockingAfterArtifactProvider(_ArtifactProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            async for event in super()._stream(call_number):
                yield event
            return
        await asyncio.Event().wait()


class _PostPublishToolLoopProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"
        self.tools_seen: list[bool] = []

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        self.tools_seen.append(bool(tools))
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderText(text="Preparing your presentation.")
            yield ProviderToolUseStart(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
                arguments={"path": "report.pptx"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderToolUseStart(tool_use_id="qa-1", tool_name="qa_check")
        yield ProviderToolUseEnd(
            tool_use_id="qa-1",
            tool_name="qa_check",
            arguments={"path": "report.pptx"},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _SelectorClone:
    current_config = SimpleNamespace(model="test/model")

    def __init__(self, provider: _ArtifactProvider) -> None:
        self.provider = provider

    def override_model(self, model: str) -> None:
        self.current_config = SimpleNamespace(model=model)
        self.provider.model = model

    def resolve(self) -> _ArtifactProvider:
        return self.provider


class _ProviderSelector:
    def __init__(self, provider: _ArtifactProvider) -> None:
        self.provider = provider

    def clone(self) -> _SelectorClone:
        return _SelectorClone(self.provider)


class _FailedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
                arguments={"path": "missing-report.pptx"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Report file is ready for download.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _RetryPublishProvider(_FailedPublishProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        yield ProviderText(text="Regenerating the presentation. ")
        yield ProviderToolUseStart(
            tool_use_id=f"publish-{call_number}",
            tool_name="publish_artifact",
        )
        yield ProviderToolUseEnd(
            tool_use_id=f"publish-{call_number}",
            tool_name="publish_artifact",
            arguments={"path": "report.pptx"},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)


class _FailedCreatePptxProvider(_FailedPublishProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="create-1", tool_name="create_pptx")
            yield ProviderToolUseEnd(
                tool_use_id="create-1",
                tool_name="create_pptx",
                arguments={"name": "report.pptx", "slides": [{"title": "Report"}]},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Report file is ready for download.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)


class _OmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "manual-big-write.html",
                    "content": "<!doctype html><title>Manual</title>",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created manual-big-write.html for you.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _OmittedInvalidPptxProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "broken.pptx",
                    "content": "this is not a PowerPoint package",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created broken.pptx for you.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _OmittedPatchPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="patch-1",
                tool_name="apply_patch",
            )
            yield ProviderToolUseEnd(
                tool_use_id="patch-1",
                tool_name="apply_patch",
                arguments={
                    "patch": (
                        "*** Begin Patch\n"
                        "*** Add File: patched.html\n"
                        "+<!doctype html><title>Patched</title>\n"
                        "*** End Patch\n"
                    ),
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created patched.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _EditedConfigProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="edit-1",
                tool_name="edit_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="edit-1",
                tool_name="edit_file",
                arguments={
                    "path": "config.json",
                    "old_text": "\"enabled\": false",
                    "new_text": "\"enabled\": true",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Updated config.json.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _MixedSizeOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            for index, payload in enumerate(
                (
                    {"path": "small.html", "content": "<title>ok</title>"},
                    {"path": "large.html", "content": "<title>" + ("x" * 80) + "</title>"},
                ),
                start=1,
            ):
                yield ProviderToolUseStart(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                )
                yield ProviderToolUseEnd(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                    arguments=payload,
                )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created small.html and large.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _MemoryJsonWriteProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "memory/cache.json",
                    "content": "{\"state\":\"internal\"}",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Updated memory/cache.json.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _SameContentOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            for index, path in enumerate(("first.html", "second.html"), start=1):
                yield ProviderToolUseStart(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                )
                yield ProviderToolUseEnd(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                    arguments={
                        "path": path,
                        "content": "<!doctype html><title>Same</title>",
                    },
                )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created first.html and second.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _PartialOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="make-1", tool_name="make_file")
            yield ProviderToolUseEnd(
                tool_use_id="make-1",
                tool_name="make_file",
                arguments={},
            )
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "second.html",
                    "content": "<!doctype html><title>Second</title>",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created runtime.txt and second.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


def _registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def make_file() -> str:
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-runtime",
                "kind": "artifact_ref",
                "name": "runtime.txt",
                "mime": "text/plain",
                "size": 4,
                "sha256": "b" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "make_file",
                "created_at": "2026-05-06T12:00:00Z",
                "download_url": (
                    "/api/v1/artifacts/art-runtime"
                    "?sessionKey=agent%3Amain%3Awebchat%3Aartifact-runtime"
                ),
            }
        )
        return "published"

    registry.register(
        ToolSpec(name="make_file", description="Make a file", parameters={}),
        make_file,
    )
    return registry


def _registry_with_write_file() -> ToolRegistry:
    registry = _registry()
    write_file = filesystem.write_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="write_file",
            description="Write a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["path", "content"],
            },
        ),
        write_file,
    )
    return registry


def _apply_patch_registry() -> ToolRegistry:
    registry = ToolRegistry()
    apply_patch = patch_tools.apply_patch.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="apply_patch",
            description="Apply a patch",
            parameters={
                "type": "object",
                "properties": {
                    "patch": {"type": "string"},
                },
                "required": ["patch"],
            },
        ),
        apply_patch,
    )
    return registry


def _edit_file_registry() -> ToolRegistry:
    registry = ToolRegistry()
    edit_file = filesystem.edit_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="edit_file",
            description="Edit a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "old_text": {"type": "string"},
                    "new_text": {"type": "string"},
                },
                "required": ["path", "old_text", "new_text"],
            },
        ),
        edit_file,
    )
    return registry


def _write_file_registry() -> ToolRegistry:
    registry = ToolRegistry()
    write_file = filesystem.write_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="write_file",
            description="Write a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["path", "content"],
            },
        ),
        write_file,
    )
    return registry


def _failed_publish_registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def publish_artifact(path: str) -> str:
        raise ToolError(f"artifact file not found: {path}")

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    return registry


def _retry_publish_registry() -> ToolRegistry:
    registry = ToolRegistry()
    calls = 0

    async def publish_artifact(path: str) -> str:
        nonlocal calls
        calls += 1
        if calls == 1:
            raise RetryableToolInputError("Regenerate the invalid PPTX and try again.")
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-retried",
                "kind": "artifact_ref",
                "name": path,
                "mime": (
                    "application/vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                "size": 8,
                "sha256": "d" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "publish_artifact",
                "created_at": "2026-07-20T00:00:00Z",
                "download_url": "/api/v1/artifacts/art-retried",
            }
        )
        return json.dumps({"status": "published", "artifact": {"name": path}})

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    return registry


def _failed_create_pptx_registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def create_pptx(slides: list[dict[str, Any]], name: str | None = None) -> str:
        raise RetryableToolInputError("The PPTX was not attached; regenerate it.")

    registry.register(
        ToolSpec(
            name="create_pptx",
            description="Create a presentation",
            parameters={
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "slides": {"type": "array", "items": {"type": "object"}},
                },
                "required": ["slides"],
            },
        ),
        create_pptx,
    )
    return registry


def _publish_then_forbidden_tool_registry() -> tuple[ToolRegistry, list[str]]:
    registry = ToolRegistry()
    forbidden_calls: list[str] = []

    async def publish_artifact(path: str) -> str:
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-published",
                "kind": "artifact_ref",
                "name": path,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "size": 8,
                "sha256": "c" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "publish_artifact",
                "created_at": "2026-05-06T12:00:00Z",
                "download_url": "/api/v1/artifacts/art-published",
            }
        )
        return json.dumps({"status": "published", "artifact": {"name": path}})

    async def qa_check(path: str) -> str:
        forbidden_calls.append(path)
        return "qa done"

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    registry.register(
        ToolSpec(
            name="qa_check",
            description="QA check",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        qa_check,
    )
    return registry, forbidden_calls


@pytest.mark.asyncio
async def test_turn_runner_streams_artifact_event_and_persists_history(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-runtime"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_ArtifactProvider()),
        tool_registry=_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make it",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]
        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert len(artifact_events) == 1
        assert artifact_events[0].id == "art-runtime"
        assert artifact_events[0].session_id == session.session_id
        assert artifact_events[0].session_key == ""
        assert artifact_events[0].download_url == "/api/v1/artifacts/art-runtime"

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["text"] == "done"
        assert payload["artifacts"][0]["id"] == "art-runtime"
        assert payload["artifacts"][0]["session_id"] == session.session_id
        assert "session_key" not in payload["artifacts"][0]
        assert "sessionKey" not in assistant.content

        class _HistoryCapture:
            def __init__(self) -> None:
                self.history = []

            def set_history(self, history) -> None:
                self.history = history

        history_capture = _HistoryCapture()
        await runner._load_history(agent=history_capture, session_key=session_key)
        assert "[generated artifact omitted: runtime.txt (text/plain)]" in str(
            history_capture.history[-1].content
        )
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_cancel_after_artifact_persists_recoverable_delivery_text(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-cancelled"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_BlockingAfterArtifactProvider()),
        tool_registry=_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )
    artifact_seen = asyncio.Event()

    async def _consume() -> None:
        async for event in runner.run(
            "make it",
            session_key,
            tool_context=tool_context,
            history_has_persisted_user=False,
            no_memory_capture=True,
        ):
            if isinstance(event, ArtifactEvent):
                artifact_seen.set()

    task = asyncio.create_task(_consume())
    try:
        await asyncio.wait_for(artifact_seen.wait(), timeout=2.0)
        await asyncio.sleep(0)
        task.cancel()
        with pytest.raises(asyncio.CancelledError):
            await task

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["id"] == "art-runtime"
        assert payload["artifacts"][0]["session_id"] == session.session_id
        assert "The generated file was delivered" in payload["text"]
        assert "[interrupted]" not in payload["text"]
    finally:
        if not task.done():
            task.cancel()
            with contextlib.suppress(asyncio.CancelledError):
                await task
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_suppresses_tools_after_successful_publish_artifact(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-stop"
    session = await manager.create(session_key)
    provider = _PostPublishToolLoopProvider()
    registry, forbidden_calls = _publish_then_forbidden_tool_registry()
    runner = TurnRunner(
        provider_selector=_ProviderSelector(provider),
        tool_registry=registry,
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make ppt",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        done = next(event for event in events if isinstance(event, DoneEvent))
        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        tool_starts = [event for event in events if isinstance(event, ToolUseStartEvent)]

        assert provider.calls == 1
        assert provider.tools_seen == [True]
        assert forbidden_calls == []
        assert [event.tool_name for event in tool_starts] == ["publish_artifact"]
        assert artifact_events[0].id == "art-published"
        assert artifact_events[0].session_id == session.session_id
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        assert "".join(text_deltas) == done.text
        assert done.text.startswith("Preparing your presentation.")
        assert "The generated file is ready" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["id"] == "art-published"
        assert "The generated file is ready" in payload["text"]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_deliverable_file_when_model_omits_publish(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-omitted"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert len(artifact_events) == 1
        assert artifact_events[0].name == "manual-big-write.html"
        assert artifact_events[0].mime == "text/html"
        assert artifact_events[0].session_id == session.session_id
        assert artifact_events[0].download_url == (
            f"/api/v1/artifacts/{artifact_events[0].id}"
        )

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["text"] == "Created manual-big-write.html for you."
        assert payload["artifacts"][0]["name"] == "manual-big-write.html"
        assert payload["artifacts"][0]["source"] == "auto_publish_omitted"
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_repeated_cancel_persists_completed_artifact_and_interrupted_transcript(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-repeated-cancel"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    publish_started = threading.Event()
    release_publish = threading.Event()
    publish_finished = threading.Event()
    real_publish = runner._stream_consumer_stage._done_handler.run_publish

    def gated_publish(inp: Any, accumulated_text: str) -> Any:
        publish_started.set()
        assert release_publish.wait(timeout=5.0), "publish was never released"
        result = real_publish(inp, accumulated_text)
        publish_finished.set()
        return result

    monkeypatch.setattr(
        runner._stream_consumer_stage._done_handler,
        "run_publish",
        gated_publish,
    )

    assistant_append_started = asyncio.Event()
    release_assistant_append = asyncio.Event()
    real_append_message = manager.append_message

    async def gated_append_message(session_key: str, **kwargs: Any) -> Any:
        if kwargs.get("role") == "assistant":
            assistant_append_started.set()
            await release_assistant_append.wait()
        return await real_append_message(session_key, **kwargs)

    monkeypatch.setattr(manager, "append_message", gated_append_message)

    async def consume() -> None:
        async for _event in runner.run(
            "make an html page",
            session_key,
            tool_context=tool_context,
            history_has_persisted_user=False,
            no_memory_capture=True,
        ):
            pass

    task = asyncio.create_task(consume())
    try:
        assert await asyncio.to_thread(publish_started.wait, 5.0)
        task.cancel("cancel-during-publish")
        release_publish.set()

        await asyncio.wait_for(assistant_append_started.wait(), timeout=5.0)
        assert publish_finished.is_set()
        task.cancel("cancel-during-transcript")
        await asyncio.sleep(0)

        release_assistant_append.set()
        with pytest.raises(asyncio.CancelledError) as exc_info:
            await task
        assert exc_info.value.args == ("cancel-during-publish",)

        transcript = await manager.get_transcript(session_key)
        assistants = [entry for entry in transcript if entry.role == "assistant"]
        assert len(assistants) == 1
        payload = json.loads(assistants[0].content)
        assert "The generated file was delivered" in payload["text"]
        assert payload["artifacts"][0]["name"] == "manual-big-write.html"

        store = ArtifactStore(str(tmp_path / "media"))
        _, artifact_path = store.resolve_for_download(
            payload["artifacts"][0]["id"],
            session_id=session.session_id,
        )
        assert artifact_path.read_text(encoding="utf-8") == (
            "<!doctype html><title>Manual</title>"
        )
    finally:
        release_publish.set()
        release_assistant_append.set()
        if not task.done():
            task.cancel()
            with contextlib.suppress(asyncio.CancelledError):
                await task
        await storage.close()


def test_auto_publish_validates_and_publishes_pptx_from_same_bytes(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import opensquilla.engine.artifact_delivery as artifact_delivery

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "brief.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(target)
    expected = target.read_bytes()
    replacement = Presentation()
    replacement.slides.add_slide(replacement.slide_layouts[5])
    replacement_output = BytesIO()
    replacement.save(replacement_output)
    replacement_payload = replacement_output.getvalue()
    validate = artifact_delivery.validate_artifact_for_delivery

    def validate_then_replace(*args: object, **kwargs: object) -> object:
        report = validate(*args, **kwargs)
        target.write_bytes(replacement_payload)
        return report

    monkeypatch.setattr(
        artifact_delivery,
        "validate_artifact_for_delivery",
        validate_then_replace,
    )

    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-valid-pptx",
        session_key="agent:main:webchat:valid-pptx",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created brief.pptx for you.",
    )

    assert result.failure_summaries == []
    assert len(result.artifacts) == 1
    artifact = result.artifacts[0]
    assert artifact["name"] == "brief.pptx"
    store = ArtifactStore(str(tmp_path / "media"))
    _, material_path = store.resolve_for_download(
        str(artifact["id"]),
        session_id="session-valid-pptx",
    )
    material = material_path.read_bytes()
    assert target.read_bytes() == replacement_payload
    assert material == expected
    Presentation(BytesIO(material))


def test_auto_publish_validates_invalid_pptx_before_persisted_dedupe(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "brief.pptx"
    invalid_payload = b"not an OOXML package"
    target.write_bytes(invalid_payload)
    media_root = tmp_path / "media"
    store = ArtifactStore(media_root)
    historical = store.publish_bytes(
        invalid_payload,
        session_id="session-invalid-pptx",
        session_key="agent:main:webchat:invalid-pptx",
        name=target.name,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="legacy",
    )
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-invalid-pptx",
        session_key="agent:main:webchat:invalid-pptx",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created brief.pptx for you.",
    )

    assert result.artifacts == []
    assert len(result.failure_summaries) == 1
    assert ctx.published_artifacts == []
    assert store.path_for(historical).read_bytes() == invalid_payload


def test_auto_publish_known_valid_pptx_reports_exact_resolved_target(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    reports = workspace / "reports"
    reports.mkdir(parents=True)
    target = reports / "brief.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(target)
    target_sha256 = hashlib.sha256(target.read_bytes()).hexdigest()
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-known-pptx",
        session_key="agent:main:webchat:known-pptx",
    )
    ctx.published_artifacts.append(
        {
            "id": "already-present",
            "sha256": target_sha256,
            "name": target.name,
        }
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": "reports/brief.pptx",
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created reports/brief.pptx for you.",
    )

    assert result.artifacts == []
    assert result.failure_summaries == []
    assert set(result.resolved_target_keys) == {
        "path:" + os.path.normcase(os.path.normpath(str(target.resolve()))),
        "name:brief.pptx",
    }


def test_auto_publish_nested_target_does_not_report_basename_as_resolved(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    reports = workspace / "reports"
    reports.mkdir(parents=True)
    target = reports / "deck.html"
    target.write_text("<title>Deck</title>", encoding="utf-8")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-nested-target",
        session_key="agent:main:webchat:nested-target",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": "reports/deck.html",
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created reports/deck.html for you.",
    )

    assert len(result.artifacts) == 1
    assert set(result.resolved_target_keys) == {
        "path:" + os.path.normcase(os.path.normpath(str(target.resolve()))),
        "name:deck.html",
    }
    assert "path:" + os.path.normcase(os.path.normpath("deck.html")) not in (
        result.resolved_target_keys
    )


def test_artifact_delivery_target_key_preserves_whitespace_and_tolerates_nul(
    tmp_path,
) -> None:
    workspace = tmp_path / "workspace"

    plain = artifact_delivery_publish_target_key(
        "deck.pptx",
        workspace_dir=workspace,
    )
    leading_space = artifact_delivery_publish_target_key(
        " deck.pptx",
        workspace_dir=workspace,
    )

    assert plain is not None
    assert leading_space is not None
    assert plain != leading_space
    assert artifact_delivery_publish_target_key(
        "bad\x00deck.pptx",
        workspace_dir=workspace,
    ) is None


@pytest.mark.parametrize(
    "target_name",
    [
        pytest.param(
            "deck:unsafe.html",
            marks=pytest.mark.skipif(
                os.name == "nt",
                reason="colon is not a legal Windows workspace filename",
            ),
        ),
        ("x" * 156) + ".html",
    ],
)
def test_auto_publish_dedupes_current_artifact_by_store_safe_name(
    tmp_path_factory: pytest.TempPathFactory,
    target_name: str,
) -> None:
    root = tmp_path_factory.mktemp("a")
    workspace = root / "w"
    workspace.mkdir()
    target = workspace / target_name
    payload = b"<title>Already published</title>"
    target.write_bytes(payload)
    media_root = root / "m"
    store = ArtifactStore(media_root)
    existing = store.publish_bytes(
        payload,
        session_id="session-safe-name",
        session_key="agent:main:webchat:safe-name",
        name=target.name,
        mime="text/html",
        source="publish_artifact",
    )
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-safe-name",
        session_key="agent:main:webchat:safe-name",
        published_artifacts=[artifact_payload(existing)],
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text=f"Created {target.name} for you.",
    )

    assert result.artifacts == []
    assert result.failure_summaries == []
    assert len(ctx.published_artifacts) == 1


def test_auto_publish_oversized_pptx_preflights_before_read_or_validation(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import opensquilla.engine.artifact_delivery as artifact_delivery

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "large.pptx"
    target.write_bytes(b"oversized")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-large-pptx",
        session_key="agent:main:webchat:large-pptx",
        artifact_max_bytes=4,
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    def unexpected_call(*args: object, **kwargs: object) -> None:
        pytest.fail("oversized PPTX must be rejected before reading or validation")

    monkeypatch.setattr(artifact_delivery.Path, "read_bytes", unexpected_call)
    monkeypatch.setattr(
        artifact_delivery,
        "validate_artifact_for_delivery",
        unexpected_call,
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created large.pptx for you.",
    )

    assert result.artifacts == []
    assert len(result.failure_summaries) == 1
    assert "per-file budget" in result.failure_summaries[0]
    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_turn_runner_rejects_invalid_omitted_pptx_and_marks_delivery_failure(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-invalid-pptx"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedInvalidPptxProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make a presentation",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        assert [event for event in events if isinstance(event, ArtifactEvent)] == []
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert "File delivery failed:" in done.text
        assert "correct or regenerate it" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_does_not_auto_publish_edited_config_json(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-edit-config"
    await manager.create(session_key)
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "config.json").write_text("{\"enabled\": false}\n", encoding="utf-8")
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_EditedConfigProvider()),
        tool_registry=_edit_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        allowed_tools={"edit_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "update config",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert assistant.content == "Updated config.json."
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_deliverable_file_created_by_apply_patch(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-patch-omitted"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPatchPublishProvider()),
        tool_registry=_apply_patch_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"apply_patch"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page with a patch",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["patched.html"]

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["name"] == "patched.html"
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_partial_omitted_artifact_delivery_failure(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-partial-failure"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_MixedSizeOmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(
                media_root=str(tmp_path / "media"),
                artifact_max_bytes=40,
            ),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make two html pages",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["small.html"]
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text
        assert "some generated files were attached" in done.text
        assert "correct or regenerate it" in done.text
        assert "no downloadable file was attached" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["name"] == "small.html"
        assert "File delivery failed:" in payload["text"]
        assert "some generated files were attached" in payload["text"]
        assert "no downloadable file was attached" not in payload["text"]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_same_content_deliverables_by_name(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-same-content"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_SameContentOmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make matching html files",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["first.html", "second.html"]
        assert artifact_events[0].sha256 == artifact_events[1].sha256

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert [artifact["name"] for artifact in payload["artifacts"]] == [
            "first.html",
            "second.html",
        ]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_omitted_deliverable_after_existing_artifact(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-partial-omitted"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_PartialOmittedPublishProvider()),
        tool_registry=_registry_with_write_file(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"make_file", "write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make two files",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["runtime.txt", "second.html"]
        assert artifact_events[0].id == "art-runtime"
        assert artifact_events[1].session_id == session.session_id

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert [artifact["name"] for artifact in payload["artifacts"]] == [
            "runtime.txt",
            "second.html",
        ]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_does_not_auto_publish_memory_json_write(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-memory-json"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_MemoryJsonWriteProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "update internal memory",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert assistant.content == "Updated memory/cache.json."
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_omitted_artifact_delivery_in_final_text(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-omitted-too-large"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(
                media_root=str(tmp_path / "media"),
                artifact_max_bytes=1,
            ),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_artifact_delivery_in_final_text(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-failed"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_FailedPublishProvider()),
        tool_registry=_failed_publish_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text
        assert "Ask me to resend the file after I correct or regenerate it." in done.text
        assert "publish_artifact" not in done.text
        assert "active workspace" not in done.text
        assert "missing-report.pptx" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "Report file is ready for download." in assistant.content
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_clears_delivery_failure_after_same_target_retry_succeeds(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-retry-success"
    await manager.create(session_key)
    provider = _RetryPublishProvider()
    runner = TurnRunner(
        provider_selector=_ProviderSelector(provider),
        tool_registry=_retry_publish_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        done = next(event for event in events if isinstance(event, DoneEvent))
        artifacts = [event for event in events if isinstance(event, ArtifactEvent)]
        assert provider.calls == 2
        assert [artifact.id for artifact in artifacts] == ["art-retried"]
        assert "File delivery failed:" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" not in assistant.content
        assert "art-retried" in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_create_pptx_delivery_in_final_text(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:create-pptx-failed"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_FailedCreatePptxProvider()),
        tool_registry=_failed_create_pptx_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert [event for event in events if isinstance(event, ArtifactEvent)] == []
        assert "File delivery failed:" in done.text
        assert "correct or regenerate it" in done.text
    finally:
        await storage.close()
